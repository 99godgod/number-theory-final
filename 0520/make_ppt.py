# -*- coding: utf-8 -*-
"""
格密碼學：SVP/CVP 的困難性與求解 — 0520 投影片生成器
所有數學以 Unicode 符號渲染（python-pptx 不支援 LaTeX 渲染）
總時長：45 分鐘 (S1 8 + S2 12 + S3 15 + S4 5 + S5 5)
輸出：presentation.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# ─── 配色 ───────────────────────────────────────────
NAVY  = RGBColor(0x0D, 0x2F, 0x6B)
BLUE  = RGBColor(0x1E, 0x50, 0xA0)
LBLUE = RGBColor(0xD6, 0xE4, 0xF7)
GOLD  = RGBColor(0xB8, 0x86, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY  = RGBColor(0x55, 0x55, 0x55)
LGRAY = RGBColor(0xF2, 0xF4, 0xF8)
GREEN = RGBColor(0x1A, 0x6B, 0x3A)
LGRN  = RGBColor(0xD4, 0xED, 0xDA)
RED   = RGBColor(0xC0, 0x30, 0x30)
LRED  = RGBColor(0xF8, 0xD7, 0xDA)

ZH = "Microsoft JhengHei"
EN = "Times New Roman"

# ─── 輔助函式 ───────────────────────────────────────
def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(slide, l, t, w, h, color, line=False):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    if not line: s.line.fill.background()
    return s

def tb(slide, text, l, t, w, h, font=ZH, size=18,
       bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    b = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    b.word_wrap = True
    tf = b.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.color.rgb = color
    return b

def footer(slide):
    rect(slide, 0, 7.1, 13.33, 0.4, LGRAY)
    tb(slide, '黃崇晉 L16141149 | 數論（一）| 2026 春季 | 國立成功大學',
       0.3, 7.13, 10, 0.3, size=11, color=GRAY)

def header(slide, title, sub=None):
    rect(slide, 0, 0, 13.33, 1.15, NAVY)
    tb(slide, title, 0.4, 0.1, 12.5, 0.75,
       size=26, bold=True, color=WHITE)
    if sub:
        tb(slide, sub, 0.4, 0.72, 12.5, 0.45,
           size=13, color=LBLUE)
    rect(slide, 0, 1.15, 13.33, 0.04, GOLD)
    footer(slide)

def bullets(title, items, sub=None, note=None):
    sld = blank()
    rect(sld, 0, 0, 13.33, 7.5, WHITE)
    header(sld, title, sub)
    b = sld.shapes.add_textbox(Inches(0.45), Inches(1.30),
                                Inches(12.4), Inches(5.5))
    b.word_wrap = True
    tf = b.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        level = item[0]; text = item[1]
        col = item[2] if len(item) > 2 else None
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(5 if level == 1 else 2)
        indent = '' if level == 1 else '    '
        bullet = '▶  ' if level == 1 else '•  '
        r = p.add_run()
        r.text = indent + bullet + text
        r.font.name = ZH
        r.font.size = Pt(19 if level == 1 else 17)
        r.font.bold = (level == 1)
        if col: r.font.color.rgb = col
        elif level == 1: r.font.color.rgb = BLUE
        else: r.font.color.rgb = BLACK
    if note:
        rect(sld, 0.4, 6.48, 12.5, 0.58, LBLUE)
        tb(sld, '💡  ' + note, 0.55, 6.51, 12.2, 0.52, size=13, color=NAVY)
    return sld

def section_divider(num, title_zh, title_en, duration):
    sld = blank()
    rect(sld, 0, 0, 13.33, 7.5, NAVY)
    rect(sld, 0, 3.0, 13.33, 0.05, GOLD)
    tb(sld, f'Section {num}', 0, 1.7, 13.33, 0.7,
       size=24, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(sld, title_zh, 0, 3.3, 13.33, 1.0,
       size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(sld, title_en, 0, 4.5, 13.33, 0.6,
       size=22, color=LBLUE, align=PP_ALIGN.CENTER, font=EN)
    tb(sld, f'⏱  {duration}', 0, 5.5, 13.33, 0.5,
       size=18, color=LBLUE, align=PP_ALIGN.CENTER)
    return sld

def formula_card(slide, l, t, w, h, formula, size=22, color=NAVY):
    rect(slide, l, t, w, h, LBLUE)
    rect(slide, l, t, 0.08, h, BLUE)
    tb(slide, formula, l + 0.2, t + 0.05, w - 0.3, h - 0.1,
       font=EN, size=size, bold=True, color=color, align=PP_ALIGN.CENTER)

import os
_IMG = os.path.join(os.path.dirname(os.path.abspath(__file__)), '..', 'minkowski_viz.png')

# ════════════════════════════════════════════════════
# Slide 1 — 標題頁
# ════════════════════════════════════════════════════
s = blank()
rect(s, 0, 0, 13.33, 7.5, NAVY)
rect(s, 0, 5.0, 13.33, 0.05, GOLD)
tb(s, 'Lattice Cryptography: Hardness and Algorithms for SVP/CVP', 0, 1.5, 13.33, 1.0,
   size=36, bold=True, color=WHITE, font=EN, align=PP_ALIGN.CENTER)
tb(s, '格密碼學：SVP/CVP 的困難性與求解', 0, 2.65, 13.33, 0.8,
   size=30, bold=True, color=LBLUE, align=PP_ALIGN.CENTER)
tb(s, 'From Minkowski\'s Theorem to Post-Quantum Security', 0, 3.6, 13.33, 0.5,
   size=18, color=LBLUE, align=PP_ALIGN.CENTER, font=EN)
tb(s, '從 Minkowski 定理到後量子安全', 0, 4.1, 13.33, 0.5,
   size=16, color=LBLUE, align=PP_ALIGN.CENTER)
tb(s, '黃崇晉 (Chung-Chin Huang)  ·  L16141149', 0, 5.25, 13.33, 0.4,
   size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, '數論（一） · 2026 春季 · 國立成功大學', 0, 5.75, 13.33, 0.4,
   size=15, color=LBLUE, align=PP_ALIGN.CENTER)
tb(s, 'Number Theory (I)  ·  Spring 2026  ·  NCKU', 0, 6.15, 13.33, 0.4,
   size=13, color=LBLUE, align=PP_ALIGN.CENTER, font=EN)

# ════════════════════════════════════════════════════
# Section 1 — 動機與背景（8 分鐘）
# ════════════════════════════════════════════════════
section_divider(1, '動機與背景', 'Motivation and Background', '8 分鐘 / minutes')

# S1-2 — RSA/ECC 基本原理
bullets('RSA 與 ECC — 現今密碼學的基石', [
    (1, 'RSA：公鑰 N = pq，安全性 = 整數分解困難'),
    (2, '古典最佳：GNFS，時間 exp(Õ((log N)^(1/3)))'),
    (2, '已知 N = pq，求 p, q 困難（IFP）'),
    (1, 'ECC：ECDLP，古典最佳 Pollard rho O(√n)'),
    (2, '已知 P, Q = kP，求純量 k 困難'),
    (2, '256-bit ECC ≈ 3072-bit RSA（更省金鑰）'),
    (1, '兩者皆為現今 HTTPS、行動銀行、TLS 基礎'),
    (2, 'Google、Apple、銀行：全球部署數百億份金鑰'),
    (2, '一旦被破解 → 整個網際網路加密生態崩潰'),
], sub='Section 1 · Slide 2 / 38',
   note='RSA 於 1977 年提出，ECC 於 1985 年提出；至今仍是最廣泛部署的公鑰系統。')

# S1-3 — Shor 演算法
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, 'Shor 演算法（1994）— 量子的致命打擊', 'Section 1 · Slide 3 / 38')
tb(s, '核心觀察：IFP 與 ECDLP 皆可化為「尋找週期 (Period Finding)」',
   0.45, 1.35, 12.4, 0.5, size=17, bold=True, color=NAVY)
formula_card(s, 0.9, 1.9, 11.5, 1.0,
             'f(x) = aˣ mod N，週期 ω 滿足 aω ≡ 1 (mod N)', size=20, color=GREEN)
tb(s, 'Quantum Fourier Transform 可在多項式時間求週期 → 直接破解 RSA / ECC',
   0.45, 3.1, 12.4, 0.45, size=16, color=BLACK)
rect(s, 0.9, 3.65, 11.5, 1.8, LRED)
rect(s, 0.9, 3.65, 0.08, 1.8, RED)
tb(s, '古典 vs 量子複雜度比較', 1.1, 3.72, 11.2, 0.4,
   size=16, bold=True, color=RED)
tb(s, 'RSA-2048：古典 GNFS  ≈  exp(Õ((log N)^(1/3)))',
   1.1, 4.15, 11.2, 0.4, size=15, color=BLACK, font=EN)
tb(s, '              →  量子 Shor  ≈  Õ((log N)³)（指數 → 多項式）',
   1.1, 4.55, 11.2, 0.4, size=15, color=BLACK, font=EN)
tb(s, 'ECDLP：古典 Pollard rho  O(√n)  →  量子 Shor  Õ((log n)³)',
   1.1, 4.95, 11.2, 0.4, size=15, color=BLACK, font=EN)
tb(s, '結論：一旦容錯量子電腦問世，RSA、DH、ECC 同時失效。', 1.1, 5.35, 11.2, 0.4,
   size=16, bold=True, color=RED)
rect(s, 0.4, 6.45, 12.5, 0.6, LBLUE)
tb(s, '💡  SVP/CVP 目前量子最佳仍需 2^Θ(n)，這正是格密碼學的立足點。',
   0.55, 6.5, 12.2, 0.5, size=14, color=NAVY)

# S1-4 — Period Finding
bullets('Period Finding 尋找週期', [
    (1, '週期函數定義：f(x + ω) = f(x) 對所有 x 成立'),
    (2, 'ω 稱為函數的週期（period）'),
    (1, 'Quantum Fourier Transform 可在多項式時間求週期'),
    (2, '與古典 DFT 類比，但作用於量子疊加態'),
    (2, '這是 Shor 演算法擊垮 RSA / ECC 的核心子程序'),
    (1, 'RSA 歸約：f(x) = aˣ mod N → 週期 ω → gcd 分解 N'),
    (2, 'N | (a^ω − 1) = (a^(ω/2)+1)(a^(ω/2)−1)'),
    (2, 'gcd(a^(ω/2) ± 1, N) 即可拆出 p 與 q'),
    (1, 'ECDLP 歸約：群上的離散對數 → 週期問題'),
    (2, '在橢圓曲線群的結構上同樣可構造週期函數'),
], sub='Section 1 · Slide 4 / 38',
   note='Period finding 是「量子電腦能做、古典電腦難做」的典型任務。')

# S1-5 — 格密碼學的登場
bullets('格密碼學的登場', [
    (1, 'Ajtai 1996：worst-case → average-case 歸約'),
    (2, '格上某些困難問題即使對量子電腦也認為困難'),
    (2, '首度給出有嚴格安全歸約的公鑰密碼學候選'),
    (1, 'LWE（Regev 2005）：帶雜訊線性方程組，公鑰 O(n²)'),
    (2, '安全性歸約自格上 SVP（量子歸約）'),
    (2, 'NIST PQC 競賽 Kyber/Dilithium 的直接前身'),
    (1, 'Ring-LWE（LPR 2010）：代數結構，公鑰 O(n)'),
    (2, 'ℤⁿ → 分圓整數環，NTT 加速乘法至 O(n log n)'),
    (1, 'SVPγ 格上最短向量問題：目前量子最佳仍需 2^Θ(n)'),
    (2, '這是格密碼學量子抵抗力的幾何根源'),
], sub='Section 1 · Slide 5 / 38')

# S1-6 — Section 1 小結：問題鏈
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, 'Section 1 小結 — 問題鏈', 'Section 1 · Slide 6 / 38')
tb(s, '問題鏈 / The Problem Chain', 0.45, 1.3, 12.4, 0.45,
   size=20, bold=True, color=NAVY)
chain = [
    ('RSA / ECC',        '整數分解 / ECDLP',             LRED),
    ('Shor 破解',        '量子多項式時間，兩者同時失效',   LRED),
    ('Lattice / SVP',   '量子難解，2^Θ(n)',               LGRN),
    ('LWE',             '帶雜訊方程組，公鑰 O(n²)',        LBLUE),
    ('Ring-LWE',        '代數結構，公鑰 O(n)，NIST 標準', LGRN),
]
for i, (k, v, c) in enumerate(chain):
    y = 1.9 + i * 0.83
    rect(s, 0.9, y, 3.5, 0.65, c)
    tb(s, k, 1.05, y + 0.12, 3.2, 0.4, size=18, bold=True, color=NAVY)
    if i < 4:
        tb(s, '→', 4.55, y + 0.14, 0.5, 0.4, size=22, bold=True, color=GOLD,
           align=PP_ALIGN.CENTER)
    rect(s, 5.1, y, 7.3, 0.65, LGRAY)
    tb(s, v, 5.3, y + 0.12, 7.0, 0.4, size=16, color=BLACK)
tb(s, '接下來：深入格的幾何結構，理解 SVP/CVP 為何困難，以及如何求解。',
   0.45, 6.42, 12.4, 0.45, size=15, bold=True, color=GREEN)

# ════════════════════════════════════════════════════
# Section 2 — 格與困難問題（12 分鐘）
# ════════════════════════════════════════════════════
section_divider(2, '格與困難問題', 'Lattices and Hard Problems', '12 分鐘 / minutes')

# S2-2 — 格的定義
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, '格 (Lattice) 的定義', 'Section 2 · Slide 8 / 38')
tb(s, '離散加法子群 / Discrete additive subgroup of ℝⁿ',
   0.45, 1.35, 12.4, 0.45, size=18, bold=True, color=NAVY)
formula_card(s, 0.9, 1.9, 11.5, 1.1,
             '𝓛  =  ⊕ᵢ₌₁ⁿ ℤ bᵢ  ⊂  ℝⁿ ,   bᵢ ∈ ℝⁿ linearly independent',
             size=20)
items_def = [
    ('基底矩陣', 'B = [b₁ | … | bₙ]，格由 B 的整數線性組合構成'),
    ('覆積 (covolume)', 'det(𝓛) = |det B|，與基底選擇無關之不變量'),
    ('相繼極小', 'λᵢ(𝓛) = inf{r : dim span(𝓛 ∩ B(0,r)) ≥ i}'),
    ('最短向量', 'λ₁(𝓛) — 第一個相繼極小，即非零最短向量長度'),
]
for i, (k, v) in enumerate(items_def):
    y = 3.2 + i * 0.65
    rect(s, 0.9, y, 3.1, 0.5, LBLUE)
    tb(s, k, 1.05, y + 0.06, 2.8, 0.4, size=15, bold=True, color=NAVY)
    tb(s, v, 4.15, y + 0.05, 8.4, 0.45, size=15, color=BLACK)
tb(s, '💡  幾何觀點：格是 ℝⁿ 中均勻分佈的點陣；覆積 = 一個基本胞（fundamental domain）的體積。',
   0.55, 6.55, 12.2, 0.5, size=13, color=NAVY)

# S2-3 — Good Basis vs Bad Basis
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, 'Good Basis vs Bad Basis — 視覺直觀', 'Section 2 · Slide 9 / 38')
tb(s, '同一個格，不同的基底呈現截然不同的樣貌。',
   0.45, 1.35, 12.4, 0.45, size=17, bold=True, color=NAVY)
rect(s, 0.45, 1.95, 5.9, 4.55, LGRN)
rect(s, 0.45, 1.95, 0.07, 4.55, GREEN)
tb(s, 'Good Basis（近正交、長度均等）', 0.6, 2.05, 5.6, 0.45,
   size=16, bold=True, color=GREEN)
good_items = [
    'b₁ ⊥ b₂（夾角近 90°）',
    '‖b₁‖ ≈ ‖b₂‖（長度相近）',
    '最短向量明顯可見',
    'Babai rounding 直接給精確解',
    'LLL reduction 後可得此形式',
]
for i, t in enumerate(good_items):
    tb(s, '▸  ' + t, 0.65, 2.55 + i * 0.5, 5.5, 0.45, size=14, color=BLACK)
rect(s, 6.95, 1.95, 5.9, 4.55, LRED)
rect(s, 6.95, 1.95, 0.07, 4.55, RED)
tb(s, 'Bad Basis（高度歪斜、夾角小）', 7.1, 2.05, 5.6, 0.45,
   size=16, bold=True, color=RED)
bad_items = [
    'b₁, b₂ 幾乎平行（夾角小）',
    '‖b₁‖ ≫ ‖b₂‖（長度差異大）',
    '最短向量隱藏其中',
    'Babai rounding 選錯胞元',
    '直接計算 B⁻¹y 誤差極大',
]
for i, t in enumerate(bad_items):
    tb(s, '▸  ' + t, 7.15, 2.55 + i * 0.5, 5.5, 0.45, size=14, color=BLACK)
rect(s, 0.45, 6.55, 12.4, 0.55, LBLUE)
tb(s, '💡  格本身不變，但 basis 的品質決定解題難度。格密碼的安全性正是利用：攻擊者拿到的是 bad basis，',
   0.6, 6.58, 12.1, 0.48, size=13, color=NAVY)

# S2-4 — SVP
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, 'SVP — Shortest Vector Problem 最短向量問題', 'Section 2 · Slide 10 / 38')
tb(s, '定義 / Definition', 0.45, 1.35, 12.4, 0.4, size=18, bold=True, color=NAVY)
formula_card(s, 0.9, 1.85, 11.5, 1.0,
             'λ₁(Λ)  =  min  ‖x‖ ,    x ∈ Λ \\ {0}', size=22, color=GREEN)
tb(s, '直觀理解 / Intuition', 0.45, 3.05, 12.4, 0.4, size=18, bold=True, color=NAVY)
svp_items = [
    'Λ 是 ℝⁿ 中離散點集；SVP = 在這些點中找離原點最近的非零點',
    '二維且 basis 漂亮時，最短向量明顯可見（眼睛可辨識）',
    '若 basis 歪扭，最短向量仍存在但極難辨識',
    '高維度（n ≥ 50）時即使 good basis 也難以找到精確解',
]
for i, t in enumerate(svp_items):
    y = 3.55 + i * 0.52
    tb(s, '▸', 0.9, y, 0.3, 0.4, size=15, bold=True, color=BLUE)
    tb(s, t, 1.2, y, 11.2, 0.48, size=15, color=BLACK)
rect(s, 0.9, 5.72, 11.5, 0.75, LRED)
rect(s, 0.9, 5.72, 0.07, 0.75, RED)
tb(s, '決策版本 GapSVP_γ：給定 Λ，判斷 λ₁ ≤ 1 還是 λ₁ > γ（中間留空隙 gap）',
   1.1, 5.8, 11.2, 0.6, size=15, bold=True, color=RED)
rect(s, 0.4, 6.52, 12.5, 0.55, LBLUE)
tb(s, '💡  SVP 與 GapSVP_γ 在小 γ 下被認為對古典 + 量子電腦均需 2^Θ(n) 時間。',
   0.55, 6.55, 12.2, 0.48, size=13, color=NAVY)

# S2-5 — CVP
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, 'CVP — Closest Vector Problem 最近向量問題', 'Section 2 · Slide 11 / 38')
tb(s, '定義 / Definition', 0.45, 1.35, 12.4, 0.4, size=18, bold=True, color=NAVY)
formula_card(s, 0.9, 1.85, 11.5, 1.0,
             'dist(y, Λ)  =  min  ‖x − y‖ ,    x ∈ Λ', size=22, color=GREEN)
tb(s, '直觀：給定目標點 y ∈ ℝⁿ（未必在格上），找格中距離 y 最近的點。',
   0.45, 3.05, 12.4, 0.45, size=16, color=BLACK)
rect(s, 0.9, 3.65, 5.6, 2.2, LBLUE)
rect(s, 0.9, 3.65, 0.07, 2.2, BLUE)
tb(s, 'BDD — Bounded Distance Decoding', 1.1, 3.72, 5.3, 0.45,
   size=15, bold=True, color=NAVY)
bdd_items = [
    'CVP 的重要特例',
    'dist(y, Λ) < λ₁ / 2',
    '此條件下答案唯一',
    'Babai NP 演算法精確求解',
    'LWE 解碼的幾何模型',
]
for i, t in enumerate(bdd_items):
    tb(s, '•  ' + t, 1.1, 4.2 + i * 0.32, 5.3, 0.3, size=13, color=BLACK)
rect(s, 6.8, 3.65, 5.6, 2.2, LRED)
rect(s, 6.8, 3.65, 0.07, 2.2, RED)
tb(s, 'CVP ≥ SVP 困難度', 7.0, 3.72, 5.3, 0.45,
   size=15, bold=True, color=RED)
hard_items = [
    'SVP 可規約至 CVP',
    '（因此 CVP 至少一樣難）',
    'CVP 在任意近似因子下',
    '均為 NP-hard',
    '精確 CVP 需指數時間',
]
for i, t in enumerate(hard_items):
    tb(s, '•  ' + t, 7.0, 4.2 + i * 0.32, 5.3, 0.3, size=13, color=BLACK)
rect(s, 0.4, 6.52, 12.5, 0.55, LBLUE)
tb(s, '💡  CVP 困難性 → BDD → LWE 安全性：這是 Kyber、Dilithium 等 NIST 後量子標準的幾何根源。',
   0.55, 6.55, 12.2, 0.48, size=13, color=NAVY)

# S2-6 — 為什麼 SVP/CVP 困難
bullets('為什麼 SVP/CVP 困難？', [
    (1, '基底歪斜時，四捨五入法（Babai rounding）失效'),
    (2, '最近格點不在最近的「胞元」中，rounding 選錯答案'),
    (1, '基底大小歸約問題：找 good basis 本身就很難'),
    (2, 'LLL 只保證指數近似，無法在多項式時間得到精確解'),
    (1, 'Minkowski 定理給上界，但找到實際最短向量是另一回事'),
    (2, 'vol(S) > 2ⁿ det(Λ) 保證存在，但不給出尋找演算法'),
    (1, '目前已知最佳古典與量子演算法'),
    (2, '精確 SVP/CVP（γ = 1）：古典 & 量子 均需 2^Θ(n)'),
    (2, '量子加速有限：BDGL sieving 僅達 2^(0.265n)（仍指數）'),
    (2, '多項式近似（γ = poly(n)）：LLL/BKZ 可達，但近似因子大'),
], sub='Section 2 · Slide 12 / 38',
   note='「困難」的本質：高維度使得格點呈指數多，而結構使得窮舉無效。')

# S2-7 — 複雜度景觀
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, '複雜度景觀 — Complexity Landscape', 'Section 2 · Slide 13 / 38')
tb(s, '格問題的複雜度依近似因子 γ 呈現清晰的分層結構：',
   0.45, 1.32, 12.4, 0.45, size=17, bold=True, color=NAVY)
rows = [
    ('問題',         '近似因子 γ',        '最佳已知複雜度',          '備註'),
    ('SVP exact',    'γ = 1',             '2^O(n)（Sieving）',       '古典 & 量子皆指數'),
    ('SVP approx',   'γ = 2^O(n^ε)',      'poly(n)（LLL）',          '多項式時間可解'),
    ('CVP exact',    'γ = 1',             '2^O(n)',                  '≥ SVP 困難'),
    ('BDD',          'dist < λ₁/2',       'poly(n)（Babai NP）',     'LWE 解碼核心'),
    ('GapSVP crypto','γ = poly(n)',        '假設困難（NIST 安全基礎）', 'LWE 歸約目標'),
]
col_x = [0.45, 2.6, 5.1, 8.35, 11.3]
col_w = [2.0, 2.3, 3.1, 2.8, 1.8]
for i, row in enumerate(rows):
    y = 1.9 + i * 0.75
    if i == 0:
        rect(s, 0.45, y, 12.4, 0.6, NAVY)
        for j, c in enumerate(row):
            tb(s, c, col_x[j] + 0.1, y + 0.12, col_w[j], 0.4,
               size=14, bold=True, color=WHITE)
    else:
        bg = LGRAY if i % 2 == 0 else WHITE
        rect(s, 0.45, y, 12.4, 0.65, bg)
        for j, c in enumerate(row):
            col = GREEN if 'poly' in c and j == 2 else (RED if '指數' in c or '困難' in c else BLACK)
            tb(s, c, col_x[j] + 0.1, y + 0.12, col_w[j], 0.42,
               size=13, color=col)

# S2-8 — Section 2 小結
bullets('Section 2 小結', [
    (1, '格 = 離散加法子群；覆積與相繼極小為核心不變量'),
    (2, 'det(𝓛) = |det B|（與基底無關）；λ₁ = 最短向量長度'),
    (1, 'SVP（找最短向量）與 CVP（找最近格點）是兩大核心困難問題'),
    (2, 'GapSVP_γ 為決策版本；BDD 是 CVP 的特例（dist < λ₁/2）'),
    (1, 'Good basis vs Bad basis：同一格，解題難度天壤之別'),
    (2, '格密碼安全性的本質：攻擊者只拿到 bad basis'),
    (1, '複雜度依近似因子分層：精確解指數，粗略近似多項式'),
    (2, 'SVP/CVP 困難性 → LWE/Ring-LWE 安全性（下節與第 3 節詳述）'),
], sub='Section 2 · Slide 14 / 38')

# ════════════════════════════════════════════════════
# Section 3 — 如何求解 SVP/CVP（15 分鐘）
# ════════════════════════════════════════════════════
section_divider(3, '如何求解 SVP 與 CVP', 'Solving SVP and CVP', '15 分鐘 / minutes')

# S3-2 — Gram-Schmidt 正交化
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, 'Gram-Schmidt 正交化（LLL 的基礎）', 'Section 3 · Slide 16 / 38')
tb(s, 'GSO 遞推公式 / Gram-Schmidt orthogonalization', 0.45, 1.32, 12.4, 0.45,
   size=17, bold=True, color=NAVY)
formula_card(s, 0.9, 1.88, 11.5, 1.3,
             'b*ᵢ  =  bᵢ  −  Σⱼ＜ᵢ μᵢⱼ b*ⱼ\n'
             'μᵢⱼ  =  ⟨bᵢ, b*ⱼ⟩ / ‖b*ⱼ‖²',
             size=20, color=GREEN)
gso_items = [
    'Gram-Schmidt 給出正交基 {b*ᵢ}，但不一定是格基（係數非整數）',
    '關鍵恆等式：∏ᵢ ‖b*ᵢ‖ = |det B| = det(Λ)（覆積不變）',
    '‖b*ᵢ‖ 的大小反映 bᵢ 相對於前面向量的「正交分量」長度',
    'LLL 條件就是利用 μᵢⱼ 與 ‖b*ᵢ‖ 對 basis 做「局部最佳化」',
]
for i, t in enumerate(gso_items):
    y = 3.4 + i * 0.55
    tb(s, '▸', 0.9, y, 0.3, 0.4, size=15, bold=True, color=BLUE)
    tb(s, t, 1.2, y, 11.2, 0.5, size=15, color=BLACK)
rect(s, 0.4, 6.48, 12.5, 0.58, LBLUE)
tb(s, '💡  GSO 的目的：把「歪扭 basis」的正交分量結構量化，作為 LLL 優化的指標。',
   0.55, 6.52, 12.2, 0.5, size=13, color=NAVY)

# S3-3 — LLL 演算法
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, 'LLL 演算法（Lenstra-Lenstra-Lovász, 1982）', 'Section 3 · Slide 17 / 38')
tb(s, 'LLL 使用兩個條件對 basis 做歸約：', 0.45, 1.32, 12.4, 0.45,
   size=17, bold=True, color=NAVY)
formula_card(s, 0.9, 1.88, 11.5, 1.3,
             '1. Size condition:  |μᵢⱼ| ≤ 1/2  for j < i\n'
             '2. Lovász condition:  ‖b*ₖ‖²  ≥  (δ − μ²ₖ,ₖ₋₁) ‖b*ₖ₋₁‖²,   δ ∈ (3/4, 1)',
             size=17, color=GREEN)
tb(s, '輸出保證：LLL-reduced basis 滿足', 0.45, 3.38, 12.4, 0.4,
   size=17, bold=True, color=NAVY)
formula_card(s, 0.9, 3.88, 11.5, 1.0,
             '‖b₁‖  ≤  2^((n−1)/2) · λ₁(Λ)    （指數近似）', size=21, color=RED)
lll_items = [
    '時間複雜度：poly(n, log B)，多項式時間可完成',
    '近似因子 2^(n/2)：維度 100 時約 10¹⁵ 倍，遠大於 1',
    '1982 年用於分解低指數 RSA、解 knapsack 密碼',
    '直觀：LLL 只做「局部最佳化」（每次交換相鄰對），無法保證全域最短',
]
for i, t in enumerate(lll_items):
    y = 5.05 + i * 0.47
    tb(s, '▸', 0.9, y, 0.3, 0.4, size=14, bold=True, color=BLUE)
    tb(s, t, 1.2, y, 11.2, 0.43, size=14, color=BLACK)
rect(s, 0.4, 6.85, 12.5, 0.3, LBLUE)
tb(s, '💡  δ 越接近 1 → 品質越好，但收斂越慢；δ = 3/4 為原始 LLL 設定。',
   0.55, 6.88, 12.2, 0.27, size=12, color=NAVY)

# S3-4 — LLL 的威力與限制
bullets('LLL 的威力與限制', [
    (1, '威力：多項式時間內找到近似最短向量'),
    (2, '這在 1982 年是突破性成果，此前無已知多項式時間方法'),
    (2, '應用：低指數 RSA 攻擊、knapsack 密碼破解、格密碼分析'),
    (1, '限制：近似因子 2^(n/2) 太大，不足以破解格密碼'),
    (2, '現代格密碼 n ≥ 256，近似因子 ≈ 2^128，完全無法攻擊'),
    (2, 'LLL 只做相鄰向量對的「局部交換」，無法跳脫局部極值'),
    (1, '改進方向：增大「窗口大小」做更全域的搜尋'),
    (2, '→ BKZ 演算法（Block Korkine-Zolotarev）'),
    (2, '以 β-維 SVP oracle 為子程序，做更大範圍的 reduction'),
], sub='Section 3 · Slide 18 / 38',
   note='LLL 是格密碼分析與弱參數攻擊的標準工具，但不足以破解設計良好的格密碼。')

# S3-5 — BKZ 演算法
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, 'BKZ 演算法（Block Korkine-Zolotarev）', 'Section 3 · Slide 19 / 38')
tb(s, '核心想法：用 β-維 SVP oracle 作為 subroutine，對每個 β×β 視窗做 reduction',
   0.45, 1.32, 12.4, 0.45, size=16, bold=True, color=NAVY)
formula_card(s, 0.9, 1.88, 11.5, 1.0,
             '近似因子  ≈  β^(n/(2β)) · poly(n)', size=22, color=GREEN)
bkz_items = [
    ('β = 2', 'LLL（相鄰兩向量對）'),
    ('β = n', '精確 SVP（指數時間，全域最佳）'),
    ('β = 20~40', '實務格密碼攻擊主力（BKZ-20 至 BKZ-40）'),
]
for i, (k, v) in enumerate(bkz_items):
    y = 3.1 + i * 0.52
    rect(s, 0.9, y, 2.2, 0.42, LBLUE)
    tb(s, k, 1.0, y + 0.05, 1.9, 0.32, size=16, bold=True, color=NAVY, font=EN)
    tb(s, v, 3.3, y + 0.05, 9.7, 0.38, size=15, color=BLACK)
formula_card(s, 0.9, 4.75, 11.5, 1.0,
             '時間複雜度：2^O(β log β) per tour × poly(n) tours', size=18, color=NAVY)
rect(s, 0.9, 5.95, 11.5, 0.65, LGRN)
rect(s, 0.9, 5.95, 0.07, 0.65, GREEN)
tb(s, '實務意義：β 越大 → 越接近精確 SVP → 攻擊能力越強，但計算成本指數增長。',
   1.1, 6.05, 11.2, 0.48, size=15, bold=True, color=GREEN)

# S3-6 — Sieving
bullets('Sieving — 精確求解 SVP', [
    (1, '基本想法：隨機取大量格點，利用「加減得更短向量」迭代篩選'),
    (2, '若 v₁, v₂ 是格點且 ‖v₁ ± v₂‖ < ‖v₁‖，則以 v₁ ± v₂ 取代 v₁'),
    (2, '迭代至收斂 → 得到最短向量'),
    (1, 'AKS Sieve（2002）：首個 2^O(n) 時間的精確 SVP 演算法'),
    (2, '理論突破：精確 SVP 從此有明確的指數上界'),
    (1, 'GaussSieve、ListSieve：實務更快的變體'),
    (2, 'GaussSieve：每步以高斯篩法減少候選向量數量'),
    (1, '空間複雜度：2^O(n)（儲存大量格點是主要瓶頸）'),
    (1, '量子 Sieving：BDGL 2016，時間 2^(0.265n)，仍指數'),
    (2, '量子加速有限：指數底數從 0.415 降至 0.265，但無法突破指數壁壘'),
], sub='Section 3 · Slide 20 / 38',
   note='Sieving 是目前精確求解高維 SVP 的最佳實作方法，但空間是最大障礙。')

# S3-7 — 求解策略比較表
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, '求解策略比較表', 'Section 3 · Slide 21 / 38')
rows = [
    ('演算法',      '目標',       '時間複雜度',        '近似因子',         '備註'),
    ('LLL',         'approx SVP', 'poly(n, log B)',    '2^(n/2)',          '多項式時間，1982 年突破'),
    ('BKZ-β',       'approx SVP', '2^O(β log β)',      'β^(n/(2β))',       '實務攻擊主力工具'),
    ('Sieving',     'exact SVP',  '2^O(n)',            '1（精確）',        '空間瓶頸 2^O(n)'),
    ('Babai Round', 'approx CVP', 'poly',              '指數誤差',         'good basis 才有用'),
    ('Babai NP',    'BDD (CVP)',  'poly',              '精確（BDD 條件）', 'dist < λ₁/2 精確解'),
]
col_x2 = [0.45, 1.85, 3.45, 5.7, 8.25, 11.0]
col_w2 = [1.3, 1.5, 2.1, 2.4, 2.6, 2.0]
for i, row in enumerate(rows):
    y = 1.38 + i * 0.83
    if i == 0:
        rect(s, 0.45, y, 12.4, 0.68, NAVY)
        for j, c in enumerate(row):
            tb(s, c, col_x2[j] + 0.08, y + 0.14, col_w2[j], 0.42,
               size=13, bold=True, color=WHITE)
    else:
        bg = LGRAY if i % 2 == 0 else WHITE
        rect(s, 0.45, y, 12.4, 0.73, bg)
        for j, c in enumerate(row):
            col = GREEN if '精確' in c or 'poly' in c else (RED if '指數' in c else BLACK)
            tb(s, c, col_x2[j] + 0.08, y + 0.14, col_w2[j], 0.5,
               size=12, color=col)
rect(s, 0.4, 6.7, 12.5, 0.42, LBLUE)
tb(s, '💡  BKZ-β 是格密碼攻擊的事實標準；Sieving 提供 SVP oracle；Babai NP 解 BDD（LWE 解碼）。',
   0.55, 6.73, 12.2, 0.35, size=12, color=NAVY)

# S3-8 — CVP 求解：Babai 演算法
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, 'CVP 求解 — Babai 演算法（1986）', 'Section 3 · Slide 22 / 38')
rect(s, 0.45, 1.32, 5.9, 4.9, LBLUE)
rect(s, 0.45, 1.32, 0.07, 4.9, BLUE)
tb(s, 'Babai Rounding（四捨五入法）', 0.6, 1.4, 5.6, 0.42,
   size=15, bold=True, color=NAVY)
rnd_items = [
    '計算坐標：c = B⁻¹ y',
    '對每個分量四捨五入：c̃ᵢ = round(cᵢ)',
    '格點近似：x̃ = B c̃',
    '快速（O(n²)），但近似因子差',
    '歪斜 basis 下誤差極大',
    '只適用於「接近正交」的 basis',
]
for i, t in enumerate(rnd_items):
    tb(s, '▸  ' + t, 0.6, 1.9 + i * 0.5, 5.5, 0.45, size=13, color=BLACK)
rect(s, 6.95, 1.32, 5.9, 4.9, LGRN)
rect(s, 6.95, 1.32, 0.07, 4.9, GREEN)
tb(s, 'Babai Nearest Plane（最近超平面法）', 7.1, 1.4, 5.6, 0.42,
   size=15, bold=True, color=GREEN)
np_items = [
    '遞推投影，逐維修正',
    '從第 n 維到第 1 維逐步決定',
    '每步投影至超平面，選最近一側',
    '在 good basis 下精確求解 BDD',
    'dist(y, Λ) < λ₁/2 時精確解',
    '複雜度：poly(n)，實用',
]
for i, t in enumerate(np_items):
    tb(s, '▸  ' + t, 7.1, 1.9 + i * 0.5, 5.5, 0.45, size=13, color=BLACK)
rect(s, 0.45, 6.32, 12.4, 0.62, LGRN)
rect(s, 0.45, 6.32, 0.07, 0.62, GREEN)
tb(s, 'BDD 的關鍵：當 dist(y, Λ) < λ₁/2 時，Babai Nearest Plane 給出精確 CVP 解。',
   0.6, 6.4, 12.1, 0.5, size=14, bold=True, color=GREEN)

# S3-9 — BDD → LWE
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, 'BDD → LWE 安全性連結', 'Section 3 · Slide 23 / 38')
tb(s, 'LWE 樣本的幾何解釋', 0.45, 1.32, 12.4, 0.42, size=18, bold=True, color=NAVY)
formula_card(s, 0.9, 1.85, 11.5, 1.0,
             'LWE 樣本：b = As + e（帶雜訊）≡ 找 Λ(A) 中離 b 最近的格點', size=17, color=GREEN)
tb(s, '歸約鏈 / Reduction Chain', 0.45, 3.05, 12.4, 0.42, size=18, bold=True, color=NAVY)
formula_card(s, 0.9, 3.58, 11.5, 1.0,
             'LWE  ←  BDD  ←  CVP  ←  SVP\n（左邊問題困難 ← 右邊問題更基礎困難）',
             size=17, color=NAVY)
security_items = [
    '若能解 BDD → 能解 LWE → 破解後量子密碼（Kyber、Dilithium）',
    '困難性：雜訊量 ‖e‖ < λ₁/2 使 BDD 精確可解，但找對應格基很難',
    '攻擊者能做的：用 BKZ 改善 basis，試圖讓 BDD 問題可解',
    '安全參數設計：選 n, q, χ 使得最好的 BKZ 攻擊仍需 2^128 操作',
]
for i, t in enumerate(security_items):
    y = 4.78 + i * 0.52
    tb(s, '▸', 0.9, y, 0.3, 0.4, size=15, bold=True, color=BLUE)
    tb(s, t, 1.2, y, 11.2, 0.48, size=15, color=BLACK)
rect(s, 0.4, 6.75, 12.5, 0.42, LGRN)
tb(s, '💡  SVP 困難 = LWE 安全：這是格密碼學安全性的幾何根源。',
   0.55, 6.78, 12.2, 0.35, size=13, color=GREEN)

# S3-10 — 安全參數
bullets('安全參數選擇（Concrete Security）', [
    (1, '攻擊者用 BKZ-β 破解需要 2^Ω(β) 計算'),
    (2, '選擇安全參數 = 讓最優 β 仍需超過目標安全等級的計算量'),
    (1, 'NIST 安全等級 1（≈ AES-128）：需要 2^128 次 BKZ 操作'),
    (1, 'Kyber-512：n = 256，q = 3329，攻擊難度 ≈ 2^118'),
    (2, 'LWE 維度 n = 256，Ring-LWE 環維度 = 256'),
    (1, 'Kyber-768：n = 384，更保守，攻擊難度 ≈ 2^182'),
    (2, '公鑰約 1184 bytes，密文約 1088 bytes'),
    (1, '選參數的三角平衡：n（維度）、q（模數）、χ（雜訊分布）'),
    (2, 'n 越大 → 越安全，但公鑰越大；q 越小 → 雜訊影響越大'),
    (2, '最終目標：用最小 n 達到目標安全等級'),
], sub='Section 3 · Slide 24 / 38',
   note='安全參數的選擇是理論（歸約）與實務（BKZ 估算）的交匯點。')

# S3-11 — Section 3 小結
bullets('Section 3 小結', [
    (1, 'LLL：多項式時間，近似因子 2^(n/2)，分析弱密碼與歷史攻擊'),
    (1, 'BKZ-β：實務格密碼攻擊主力，近似因子 β^(n/(2β))'),
    (2, 'β 越大越強；BKZ-20 至 BKZ-40 為當前 NIST 參數估算標準'),
    (1, 'Sieving：精確 SVP，時間空間皆 2^O(n)，量子加速至 2^(0.265n)'),
    (1, 'Babai NP：BDD 精確求解 CVP，LWE 解碼的幾何核心'),
    (2, 'dist(y, Λ) < λ₁/2 條件下，good basis 給出精確解'),
    (1, 'SVP/CVP 困難性 → LWE/Ring-LWE 安全性（歸約鏈）'),
    (2, 'NIST Kyber（FIPS 203）與 Dilithium（FIPS 204）的安全基礎'),
], sub='Section 3 · Slide 25 / 38')

# ════════════════════════════════════════════════════
# Section 4 — Minkowski 定理與幾何基礎（5 分鐘）
# ════════════════════════════════════════════════════
section_divider(4, 'Minkowski 定理與幾何基礎',
                "Minkowski's Theorem and Geometry", '5 分鐘 / minutes')

# S4-2 — Minkowski 第一定理
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, 'Minkowski 第一定理（1889）', 'Section 4 · Slide 27 / 38')
tb(s, '幾何數論的奠基定理', 0.45, 1.35, 12.4, 0.42, size=18, bold=True, color=NAVY)
rect(s, 0.9, 1.9, 11.5, 1.65, LGRN)
rect(s, 0.9, 1.9, 0.08, 1.65, GREEN)
tb(s, 'Theorem (Minkowski, 1889).', 1.1, 1.98, 11.2, 0.42,
   size=18, bold=True, color=GREEN, font=EN)
tb(s, '若 S ⊂ ℝⁿ 為中心對稱凸體，且 vol(S) > 2ⁿ · det(𝓛)，',
   1.1, 2.45, 11.2, 0.45, size=17, color=BLACK)
tb(s, '則 S 含有非零格點 v ∈ 𝓛 ∩ S \\ {0}。',
   1.1, 2.92, 11.2, 0.45, size=17, color=BLACK)
tb(s, '三個關鍵詞 / Three keywords', 0.45, 3.82, 12.4, 0.42,
   size=17, bold=True, color=NAVY)
keywords = [
    ('中心對稱', 'centrally symmetric：x ∈ S ⟹ −x ∈ S'),
    ('凸體',     'convex body：任兩點線段含於 S；S 為有界閉集'),
    ('體積條件', 'vol(S) > 2ⁿ · det(𝓛) 為臨界條件'),
]
for i, (k, v) in enumerate(keywords):
    y = 4.38 + i * 0.55
    rect(s, 0.9, y, 2.6, 0.45, LBLUE)
    tb(s, k, 1.05, y + 0.05, 2.3, 0.35, size=15, bold=True, color=NAVY)
    tb(s, v, 3.7, y + 0.05, 8.7, 0.42, size=14, color=BLACK)
rect(s, 0.4, 6.5, 12.5, 0.58, LBLUE)
tb(s, '💡  下兩張投影片將套用此定理推導 SVP 上界 λ₁(Λ) ≤ √n · det(Λ)^(1/n)。',
   0.55, 6.53, 12.2, 0.5, size=13, color=NAVY)

# S4-3 — SVP 上界推導（含視覺化圖片）
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, 'SVP 上界推導（Minkowski → λ₁ 上界）', 'Section 4 · Slide 28 / 38')
tb(s, 'Step 1：取中心對稱立方體', 0.45, 1.32, 12.4, 0.42, size=17, bold=True, color=NAVY)
formula_card(s, 0.9, 1.85, 11.5, 1.0,
             'Cᵣ := [−r, r]ⁿ,  vol(Cᵣ) = (2r)ⁿ > 2ⁿ det(Λ)  ⟺  r > det(Λ)^(1/n)',
             size=18, color=GREEN)
tb(s, 'Step 2：範數不等式', 0.45, 3.05, 12.4, 0.42, size=17, bold=True, color=NAVY)
formula_card(s, 0.9, 3.57, 11.5, 0.92,
             '‖v‖∞ ≤ r  ⟹  ‖v‖₂ ≤ √n · ‖v‖∞ ≤ √n · r',
             size=18, color=NAVY)
tb(s, '令 r ↘ det(Λ)^(1/n)，得結論：', 0.45, 4.65, 12.4, 0.4,
   size=16, color=BLACK)
rect(s, 0.9, 5.15, 11.5, 1.0, LGRN)
rect(s, 0.9, 5.15, 0.08, 1.0, GREEN)
tb(s, 'λ₁(Λ)  ≤  √n  ·  det(Λ)^(1/n)', 1.1, 5.32, 11.2, 0.75,
   size=28, bold=True, color=GREEN, font=EN, align=PP_ALIGN.CENTER)
rect(s, 0.4, 6.52, 12.5, 0.55, LBLUE)
tb(s, '💡  Hermite 常數 γₙ 給更緊版本：λ₁² ≤ γₙ det(Λ)^(2/n)，γₙ ≤ 2n/πe + o(n)。',
   0.55, 6.55, 12.2, 0.48, size=13, color=NAVY)

# S4-3b — Minkowski 視覺化
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, 'Minkowski 定理 — 視覺化 (n = 2)', 'Section 4 · Slide 29 / 38')
tb(s, '整數格 Λ = ℤ² 與中心對稱凸體（圓盤）：臨界體積 = 2² det(Λ) = 4',
   0.45, 1.32, 12.4, 0.45, size=16, bold=True, color=NAVY,
   align=PP_ALIGN.CENTER)
if os.path.exists(_IMG):
    s.shapes.add_picture(_IMG, Inches(1.665), Inches(1.85), width=Inches(10))
rect(s, 0.4, 6.6, 12.5, 0.45, LBLUE)
tb(s, '💡  (a) vol(S) ≤ 4：S 可能只含原點。  (b) vol(S) > 4：S 必含非零格點 v ∈ Λ \\ {0}。',
   0.55, 6.62, 12.2, 0.4, size=13, color=NAVY)

# S4-4 — 定理的意義
bullets('Minkowski 定理的意義', [
    (1, '連結覆積（幾何）與最短向量（算術）'),
    (2, 'det(Λ) 越大（基本胞越大）→ λ₁ 也可以更大'),
    (2, '幾何（面積/體積）與數論（格點距離）的橋樑'),
    (1, '給出格密碼安全參數的理論下界'),
    (2, 'λ₁ 越大 → 格越「稀疏」→ SVP 越難 → 密碼越安全'),
    (2, '設計格密碼時，選參數使 λ₁ 遠超攻擊演算法的求解能力'),
    (1, '反向應用：攻擊分析'),
    (2, 'Minkowski 界是 λ₁ 的理論上界，幫助估算攻擊難度'),
    (2, 'Hermite 因子衡量 LLL/BKZ 輸出品質：‖b₁‖ / det(Λ)^(1/n)'),
], sub='Section 4 · Slide 30 / 38',
   note='Minkowski 定理是格密碼安全參數設計的理論基礎，也是 Course Objective 5。')

# ════════════════════════════════════════════════════
# Section 5 — 應用與開放問題（5 分鐘）
# ════════════════════════════════════════════════════
section_divider(5, '應用與開放問題', 'Applications and Open Problems', '5 分鐘 / minutes')

# S5-2 — NIST 後量子標準
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, 'NIST 後量子標準（2024 年 8 月公布）', 'Section 5 · Slide 32 / 38')
tb(s, '歷時 8 年的全球競賽（2016–2024）',
   0.45, 1.3, 12.4, 0.42, size=15, color=GRAY)
rect(s, 0.9, 1.85, 11.5, 2.1, LGRN)
rect(s, 0.9, 1.85, 0.08, 2.1, GREEN)
tb(s, 'FIPS 203 — 金鑰交換 / Key Encapsulation (ML-KEM = Kyber)',
   1.1, 1.95, 11.2, 0.42, size=18, bold=True, color=GREEN)
tb(s, '▸ 後量子版本的 RSA / Diffie–Hellman 握手；用於建立加密通道',
   1.1, 2.42, 11.2, 0.4, size=15, color=BLACK)
tb(s, '▸ 數學基礎：Module-LWE（Ring-LWE 的模組化推廣）',
   1.1, 2.82, 11.2, 0.4, size=15, color=BLACK)
tb(s, '▸ 公鑰約 800 bytes，密文約 1 KB；速度快於 RSA-2048',
   1.1, 3.22, 11.2, 0.4, size=14, color=GRAY)
rect(s, 0.9, 4.1, 11.5, 2.1, LBLUE)
rect(s, 0.9, 4.1, 0.08, 2.1, BLUE)
tb(s, 'FIPS 204 — 數位簽章 / Digital Signature (ML-DSA = Dilithium)',
   1.1, 4.2, 11.2, 0.42, size=18, bold=True, color=BLUE)
tb(s, '▸ 後量子版本的 RSA / ECDSA 簽章；用於 TLS 憑證、軟體簽章',
   1.1, 4.65, 11.2, 0.4, size=15, color=BLACK)
tb(s, '▸ 數學基礎：Module-LWE / Module-SIS（SVP/CVP 困難性）',
   1.1, 5.05, 11.2, 0.4, size=15, color=BLACK)
tb(s, '▸ 簽章約 2.4 KB，公鑰約 1.3 KB',
   1.1, 5.45, 11.2, 0.4, size=14, color=GRAY)
rect(s, 0.4, 6.52, 12.5, 0.55, LBLUE)
tb(s, '💡  SVP/CVP 困難性是 Kyber 與 Dilithium 安全性的幾何根基，也是本報告的主題。',
   0.55, 6.55, 12.2, 0.48, size=14, bold=True, color=NAVY)

# S5-3 — 開放問題
bullets('開放問題', [
    (1, 'SVP/CVP 的量子複雜度下界是否為 2^Ω(n)？'),
    (2, '目前只有 2^O(n) 上界（Sieving），下界幾乎未知'),
    (1, 'LLL/BKZ 的精確最壞情形複雜度？'),
    (2, 'BKZ 實作行為（tour 數）遠優於理論上界，缺乏精確分析'),
    (1, 'Ideal-SVP vs SVP：理想格是否真的比任意格容易？'),
    (2, '目前無已知演算法可利用理想格結構加速，但歸約仍較弱'),
    (1, 'Classical worst-case → average-case 歸約？'),
    (2, 'Regev 歸約為量子的；古典版本至今仍是開放問題'),
    (1, '格密碼 vs 其他 PQC：NTRU、McEliece、Hash-based 的比較'),
    (2, 'McEliece：超大公鑰；Hash-based：簽章大；NTRU：參數選擇複雜'),
], sub='Section 5 · Slide 33 / 38',
   note='格密碼學在理論與實作上都有大量未解問題，是當前最活躍的研究方向之一。')

# S5-4 — 結尾
s = blank(); rect(s, 0, 0, 13.33, 7.5, NAVY)
rect(s, 0, 3.2, 13.33, 0.05, GOLD)
tb(s, 'Thank you!', 0, 1.2, 13.33, 1.0,
   size=64, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=EN)
tb(s, '感謝聆聽 · Q & A', 0, 2.35, 13.33, 0.7,
   size=32, color=LBLUE, align=PP_ALIGN.CENTER)
tb(s, '核心問題鏈 / Core Problem Chain', 0, 3.5, 13.33, 0.45,
   size=18, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
tb(s, 'SVP/CVP 困難  →  LWE/Ring-LWE  →  NIST PQC 標準（Kyber / Dilithium）',
   0, 4.0, 13.33, 0.5, size=16, color=LBLUE, align=PP_ALIGN.CENTER)
tb(s, 'Minkowski 定理  →  λ₁ 上界  →  安全參數設計的理論基礎',
   0, 4.5, 13.33, 0.5, size=15, color=LBLUE, align=PP_ALIGN.CENTER)
tb(s, '主要參考文獻 / Main References', 0.5, 5.15, 12.3, 0.4,
   size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
refs = [
    'O. Regev, "On lattices, learning with errors, random linear codes, and cryptography." J. ACM 56(6), 2009.',
    'M. Ajtai, "Generating hard instances of lattice problems." STOC 1996.',
    'D. Micciancio & S. Goldwasser, "Complexity of Lattice Problems." Kluwer, 2002.',
    'NIST FIPS 203 / 204, 2024.    ·    D. Micciancio & O. Regev, "Lattice-based cryptography." PQC, 2009.',
]
for i, r in enumerate(refs):
    tb(s, r, 0.5, 5.58 + i * 0.27, 12.3, 0.26,
       size=11, color=LBLUE, align=PP_ALIGN.CENTER, font=EN)
tb(s, '黃崇晉  ·  L16141149  ·  數論（一）  ·  2026 春季  ·  國立成功大學',
   0, 6.7, 13.33, 0.4, size=12, color=LBLUE, align=PP_ALIGN.CENTER)

# ─── Save ──────────────────────────────────────────
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'presentation.pptx')
prs.save(out)
print(f'Saved: {out}')
print(f'Total slides: {len(prs.slides)}')
