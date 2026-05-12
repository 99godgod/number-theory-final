# -*- coding: utf-8 -*-
"""
Ring-LWE 與理想格密碼學 — v2 投影片生成器
所有數學以 Unicode 符號渲染（python-pptx 不支援 LaTeX 渲染）
總時長：45 分鐘 (S1 10 + S2 10 + S3 10 + S4 10 + S5 5)
輸出：presentation_v2.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# ─── 配色 ───────────────────────────────────────────
NAVY  = RGBColor(0x0D, 0x2F, 0x6B)
BLUE  = RGBColor(0x1E, 0x50, 0xA0)
LBLUE = RGBColor(0xD6, 0xE4, 0xF7)
GOLD  = RGBColor(0xB8, 0x86, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY  = RGBColor(0x55, 0x55, 0x55)
LGRAY = RGBColor(0xF2, 0xF4, 0xF8)
GREEN = RGBColor(0x1A, 0x6B, 0x3A)
LGRN  = RGBColor(0xD4, 0xED, 0xDA)
RED   = RGBColor(0xC0, 0x30, 0x30)
LRED  = RGBColor(0xF8, 0xD7, 0xDA)

ZH = "Microsoft JhengHei"
EN = "Times New Roman"

# ─── 輔助函式 ───────────────────────────────────────
def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(slide, l, t, w, h, color, line=False):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    if not line: s.line.fill.background()
    return s

def tb(slide, text, l, t, w, h, font=ZH, size=18,
       bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    b = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    b.word_wrap = True
    tf = b.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.color.rgb = color
    return b

def footer(slide):
    rect(slide, 0, 7.1, 13.33, 0.4, LGRAY)
    tb(slide, "黃崇晉 L16141149 | 數論（一）| 2026 春季 | 國立成功大學",
       0.3, 7.13, 10, 0.3, size=11, color=GRAY)

def header(slide, title, sub=None):
    rect(slide, 0, 0, 13.33, 1.15, NAVY)
    tb(slide, title, 0.4, 0.1, 12.5, 0.75,
       size=26, bold=True, color=WHITE)
    if sub:
        tb(slide, sub, 0.4, 0.72, 12.5, 0.45,
           size=13, color=LBLUE)
    rect(slide, 0, 1.15, 13.33, 0.04, GOLD)
    footer(slide)

def bullets(title, items, sub=None, note=None):
    sld = blank()
    rect(sld, 0, 0, 13.33, 7.5, WHITE)
    header(sld, title, sub)
    b = sld.shapes.add_textbox(Inches(0.45), Inches(1.30),
                                Inches(12.4), Inches(5.5))
    b.word_wrap = True
    tf = b.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        level = item[0]; text = item[1]
        col = item[2] if len(item) > 2 else None
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(5 if level == 1 else 2)
        indent = "" if level == 1 else "    "
        bullet = "▶  " if level == 1 else "•  "
        r = p.add_run()
        r.text = indent + bullet + text
        r.font.name = ZH
        r.font.size = Pt(19 if level == 1 else 17)
        r.font.bold = (level == 1)
        if col: r.font.color.rgb = col
        elif level == 1: r.font.color.rgb = BLUE
        else: r.font.color.rgb = BLACK
    if note:
        rect(sld, 0.4, 6.48, 12.5, 0.58, LBLUE)
        tb(sld, "💡  " + note, 0.55, 6.51, 12.2, 0.52, size=13, color=NAVY)
    return sld

def section_divider(num, title_zh, title_en, duration):
    sld = blank()
    rect(sld, 0, 0, 13.33, 7.5, NAVY)
    rect(sld, 0, 3.0, 13.33, 0.05, GOLD)
    tb(sld, f"Section {num}", 0, 1.7, 13.33, 0.7,
       size=24, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(sld, title_zh, 0, 3.3, 13.33, 1.0,
       size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(sld, title_en, 0, 4.5, 13.33, 0.6,
       size=22, color=LBLUE, align=PP_ALIGN.CENTER, font=EN)
    tb(sld, f"⏱  {duration}", 0, 5.5, 13.33, 0.5,
       size=18, color=LBLUE, align=PP_ALIGN.CENTER)
    return sld

def formula_card(slide, l, t, w, h, formula, size=22, color=NAVY):
    rect(slide, l, t, w, h, LBLUE)
    rect(slide, l, t, 0.08, h, BLUE)
    tb(slide, formula, l + 0.2, t + 0.05, w - 0.3, h - 0.1,
       font=EN, size=size, bold=True, color=color, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════
# Slide 1 — Title
# ════════════════════════════════════════════════════
s = blank()
rect(s, 0, 0, 13.33, 7.5, NAVY)
rect(s, 0, 5.0, 13.33, 0.05, GOLD)
tb(s, "Ring-LWE and Ideal Lattice Cryptography", 0, 1.7, 13.33, 1.0,
   size=42, bold=True, color=WHITE, font=EN, align=PP_ALIGN.CENTER)
tb(s, "Ring-LWE 與理想格密碼學", 0, 2.85, 13.33, 0.8,
   size=32, bold=True, color=LBLUE, align=PP_ALIGN.CENTER)
tb(s, "Algebraic Number Theory Meets Post-Quantum Security", 0, 3.85, 13.33, 0.5,
   size=18, color=LBLUE, align=PP_ALIGN.CENTER, font=EN)
tb(s, "代數數論與後量子安全的交匯", 0, 4.35, 13.33, 0.5,
   size=16, color=LBLUE, align=PP_ALIGN.CENTER)
tb(s, "黃崇晉 (Chung-Chin Huang)  ·  L16141149", 0, 5.3, 13.33, 0.4,
   size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "數論（一） · 2026 春季 · 國立成功大學", 0, 5.8, 13.33, 0.4,
   size=15, color=LBLUE, align=PP_ALIGN.CENTER)
tb(s, "Number Theory (I)  ·  Spring 2026  ·  NCKU", 0, 6.2, 13.33, 0.4,
   size=13, color=LBLUE, align=PP_ALIGN.CENTER, font=EN)

# ════════════════════════════════════════════════════
# Section 1 — RSA/ECC + Motivation (10 min)
# ════════════════════════════════════════════════════
section_divider(1, "為什麼需要後量子密碼學？",
                "Why Do We Need Post-Quantum Cryptography?",
                "10 分鐘 / minutes")

# Slide 3 — RSA basics
bullets("RSA — 整數分解問題", [
    (1, "RSA (Rivest–Shamir–Adleman, 1977)"),
    (2, "公鑰：N = p · q（兩大質數乘積）、e 為公鑰指數"),
    (2, "私鑰：d 滿足  ed ≡ 1 (mod φ(N))"),
    (2, "加密  m → c = mᵉ mod N；解密  c → m = cᵈ mod N"),
    (1, "安全性建基於 Integer Factoring Problem"),
    (2, "已知 N，找 p, q 困難"),
    (2, "若能分解 N，即可從公鑰算出私鑰，整個系統崩潰"),
    (1, "古典最佳演算法:General Number Field Sieve (GNFS)"),
    (2, "次指數時間  exp(Õ((log N)^(1/3)))"),
    (2, "目前 RSA-2048 仍認為安全"),
], sub="Section 1 · Slide 3 / 41",
   note="當前 HTTPS、銀行、SSH 等基礎設施大量依賴 RSA。")

# Slide 4 — ECC basics
bullets("ECC — 橢圓曲線離散對數問題", [
    (1, "Elliptic Curve Cryptography (Miller / Koblitz, 1985)"),
    (2, "在橢圓曲線  E: y² = x³ + ax + b  上的點集成阿貝爾群"),
    (2, "公鑰：Q = kP（純量乘法）、P 為公開生成元"),
    (2, "私鑰：純量 k"),
    (1, "安全性:Elliptic Curve Discrete Logarithm Problem (ECDLP)"),
    (2, "已知 P, Q = kP，求 k 困難"),
    (2, "比 RSA 更省金鑰：256-bit ECC ≈ 3072-bit RSA"),
    (1, "古典最佳演算法:Pollard rho"),
    (2, "O(√n) 時間（n 為群階）"),
    (2, "256-bit 曲線目前仍認為安全"),
], sub="Section 1 · Slide 4 / 41",
   note="ECDSA 用於 TLS 憑證、區塊鏈簽章；NIST P-256、Curve25519 廣泛部署。")


# Slide 5 — Period Finding
bullets("Period Finding 尋找週期", [
    (1, "週期函數定義 / Periodic function"),
    (2, "f : (x₁, …, xₙ) ↦ f(x₁, …, xₙ) 為 (ω₁, …, ωₙ)-週期，若"),
    (2, "f(x₁ + ω₁, …, xₙ + ωₙ) = f(x₁, …, xₙ)  對所有 (x₁, …, xₙ) ∈ dom(f) 成立"),
    (1, "為什麼重要？"),
    (2, "量子電腦擅長尋找週期函數的週期 (period finding)"),
    (2, "Quantum Fourier Transform (QFT) 可在多項式時間求週期"),
    (2, "Shor 演算法的核心子程序"),
], sub="Section 1 · Slide 5 / 41",
   note="Period finding 是 Shor 演算法擊垮 RSA / ECC 的「關鍵齒輪」。")

# Slide 6 — Integer Factoring
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, "Integer Factoring 質因數分解", "Section 1 · Slide 6 / 41")
tb(s, "[Factorization Problem] 給定 RSA 模數 N = p · q，求質數 p, q。",
   0.45, 1.35, 12.4, 0.5, size=17, bold=True, color=NAVY)
items_if = [
    "選取隨機整數 a ∈ ℤ_N（不失一般性，假設 gcd(a, N) = 1；否則已得 N 的因數）",
    "考慮單變量函數  f : x ↦ f(x) = aˣ mod N",
    "Period finding 子程序找到 ω 使 f(x + ω) = f(x)，即 a^ω ≡ 1 (mod N)",
    "基於 Miller 方法的演算法可由此分解 N：",
]
for i, t in enumerate(items_if):
    y = 2.0 + i * 0.55
    tb(s, "▸", 0.9, y, 0.3, 0.4, size=16, bold=True, color=BLUE)
    tb(s, t, 1.2, y, 11.2, 0.55, size=15, color=BLACK)
formula_card(s, 0.9, 4.3, 11.5, 1.2,
             "N | (a^ω − 1)   ⟹   p · q | (a^(ω/2) + 1)(a^(ω/2) − 1)",
             size=22, color=GREEN)
tb(s, "再用 gcd(a^(ω/2) ± 1, N) 即可拆出 p 與 q（多數情況成功）。",
   0.45, 5.65, 12.4, 0.5, size=16, color=BLACK)
rect(s, 0.4, 6.45, 12.5, 0.6, LBLUE)
tb(s, "💡  Integer Factoring 經 period finding 歸約後，量子可在 Õ((log N)³) 時間破解。",
   0.55, 6.5, 12.2, 0.5, size=14, color=NAVY)

# Slide 7 — Shor's algorithm
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, "Shor 演算法 (1994) — 致命的量子打擊", "Section 1 · Slide 7 / 41")
tb(s, "Peter W. Shor, “Algorithms for Quantum Computation: Discrete Logarithms and Factoring”, FOCS 1994.",
   0.45, 1.32, 12.4, 0.4, size=15, color=GRAY, font=EN)
formula_card(s, 0.9, 1.85, 11.5, 1.0,
             "Shor:   IFP, ECDLP   ∈   BQP   (quantum polynomial time)", size=22)
tb(s, "核心觀察：兩個問題皆可化為求週期 (period finding)，量子傅立葉轉換以多項式時間求出。",
   0.45, 3.0, 12.4, 0.6, size=18, color=BLACK)
rect(s, 0.9, 3.7, 11.5, 1.7, LRED)
rect(s, 0.9, 3.7, 0.08, 1.7, RED)
tb(s, "古典 vs 量子複雜度比較   /   Classical vs Quantum",
   1.1, 3.78, 11.2, 0.4, size=16, bold=True, color=RED)
tb(s, "整數分解 (RSA-2048):   classical GNFS  ≈  exp(Õ((log N)^(1/3)))",
   1.1, 4.18, 11.2, 0.4, size=15, color=BLACK, font=EN)
tb(s, "                                      →   quantum Shor   ≈   Õ((log N)³)",
   1.1, 4.45, 11.2, 0.4, size=15, color=BLACK, font=EN)
tb(s, "ECDLP:   classical Pollard rho  ≈  O(√n)    →    quantum Shor  ≈  Õ((log n)³)",
   1.1, 4.78, 11.2, 0.4, size=15, color=BLACK, font=EN)
tb(s, "結論：一旦容錯量子電腦問世，RSA、Diffie–Hellman、ECC 將同時失效。",
   1.1, 5.08, 11.2, 0.4, size=16, bold=True, color=RED)
tb(s, "💡  並非所有量子演算法都這麼快，例如 SVP/CVP 目前仍只有指數時間量子演算法。",
   0.55, 6.55, 12.2, 0.5, size=14, color=NAVY)

# Slide 6 — Harvest now decrypt later
bullets("威脅模型 — Harvest Now, Decrypt Later", [
    (1, "迫切性：威脅不需要等到量子電腦造好"),
    (2, "今日攔截的密文，可被攻擊者保存"),
    (2, "未來量子電腦上線後再解密 → 「先攔截、後解密」"),
    (1, "對長期機密影響重大"),
    (2, "醫療紀錄、外交電文、商業機密、身分認證金鑰"),
    (2, "若加密生命週期 > 量子電腦問世時間 → 已經暴露"),
    (1, "估計時程"),
    (2, "NIST 預估 2030–2035 量子電腦可能達破解規模"),
    (2, "PQC 標準必須提前部署 — 這是現在進行式"),
], sub="Section 1 · Slide 8 / 41",
   note="“Y2Q” = Years to Quantum；產業界已開始將 PQC 並列部署於現有系統。")

# Slide 7 — Lattice-based crypto intro
bullets("解方:格密碼學 (Lattice-based Cryptography)", [
    (1, "Ajtai 1996 — 第一個 worst-case 到 average-case 的歸約"),
    (2, "格上某些「困難問題」即使對量子電腦也認為困難"),
    (2, "SVP_γ（最短向量問題）：2^Θ(n) 時間"),
    (1, "Regev 2005 — Learning With Errors (LWE)"),
    (2, "公鑰加密第一個有效 PQC 候選"),
    (2, "缺點：公鑰大小  O(n²)"),
    (1, "Lyubashevsky–Peikert–Regev 2010 — Ring-LWE"),
    (2, "結合代數數論：ℤⁿ → 𝒪_K"),
    (2, "公鑰縮為 O(n)、乘法 O(n log n)（NTT）"),
], sub="Section 1 · Slide 9 / 41",
   note="本報告焦點：以代數數論的觀點，理解 Ring-LWE 為何兼具安全性與效率。")

# Slide 8 — Section 1 summary / bridge
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, "Section 1 小結 — 通往 Ring-LWE 的脈絡", "Section 1 · Slide 10 / 41")
tb(s, "問題鏈 / The Problem Chain", 0.45, 1.4, 12.4, 0.5,
   size=20, bold=True, color=NAVY)
items = [
    ("RSA / ECC", "整數分解 / ECDLP", LRED),
    ("Shor (1994)", "量子多項式破解兩者", LRED),
    ("Lattice", "SVP_γ 量子難解", LGRN),
    ("LWE (2005)", "公鑰大  O(n²)", LBLUE),
    ("Ring-LWE (2010)", "公鑰小 O(n)、NTT 加速", LGRN),
]
for i, (k, v, c) in enumerate(items):
    y = 2.0 + i * 0.8
    rect(s, 0.9, y, 4.0, 0.65, c)
    tb(s, k, 1.1, y + 0.12, 3.6, 0.4, size=18, bold=True, color=NAVY)
    rect(s, 5.1, y, 7.3, 0.65, LGRAY)
    tb(s, v, 5.3, y + 0.12, 7.0, 0.4, size=16, color=BLACK)
tb(s, "接下來：Section 2 — 進入格與幾何數論的世界，理解為什麼 SVP 困難。",
   0.45, 6.4, 12.4, 0.5, size=15, bold=True, color=GREEN)

# ════════════════════════════════════════════════════
# Section 2 — Lattices and Geometry of Numbers (10 min)
# ════════════════════════════════════════════════════
section_divider(2, "格與幾何數論",
                "Lattices and Geometry of Numbers",
                "10 分鐘 / minutes")

# Slide 10 — Lattice definition
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, "格 (Lattice) 的定義", "Section 2 · Slide 12 / 41")
tb(s, "離散加法子群 / Discrete additive subgroup", 0.45, 1.4, 12.4, 0.5,
   size=18, bold=True, color=NAVY)
formula_card(s, 0.9, 2.0, 11.5, 1.1,
             "𝓛  =  ⊕ᵢ₌₁ⁿ ℤ bᵢ  ⊂  ℝⁿ ,   bᵢ ∈ ℝⁿ linearly independent",
             size=20)
items_zh = [
    ("基底矩陣", "B = [b₁ | … | bₙ]"),
    ("覆積 (covolume)", "det(𝓛) = |det B|，與基底選擇無關之不變量"),
    ("相繼極小", "λᵢ(𝓛) := inf{r : dim span(𝓛 ∩ B(0,r)) ≥ i}"),
    ("最短向量", "λ₁(𝓛) — 第一個相繼極小，即非零最短向量長度"),
]
for i, (k, v) in enumerate(items_zh):
    y = 3.3 + i * 0.65
    rect(s, 0.9, y, 3.0, 0.5, LBLUE)
    tb(s, k, 1.05, y + 0.06, 2.7, 0.4, size=16, bold=True, color=NAVY)
    tb(s, v, 4.1, y + 0.05, 8.5, 0.45, size=15, color=BLACK)
tb(s, "💡  幾何觀點：格是 ℝⁿ 中均勻分佈的點陣，覆積 = 一個基本胞的體積。",
   0.55, 6.55, 12.2, 0.5, size=14, color=NAVY)

# Slide 11 — Hard problems
bullets("格上的困難問題 (Hard Problems on Lattices)", [
    (1, "SVP_γ — Shortest Vector Problem"),
    (2, "找  v ∈ 𝓛 \\ {0}  使  ‖v‖ ≤ γ · λ₁(𝓛)"),
    (2, "最自然的格問題；γ = 1 時即「找最短向量」"),
    (1, "GapSVP_γ — 決策版本"),
    (2, "判斷  λ₁ ≤ 1  或  λ₁ > γ"),
    (1, "BDD_d — Bounded Distance Decoding"),
    (2, "給定  t ∈ ℝⁿ，dist(t, 𝓛) ≤ d · λ₁ / 2"),
    (2, "找最近的格點（解碼）"),
    (1, "目前最佳演算法（古典與量子）"),
    (2, "在小 γ 下需要  2^Θ(n)  時間"),
    (2, "BKZ + Sieving 為實作主流，仍指數複雜度"),
], sub="Section 2 · Slide 13 / 41")


# Slide 14 — SVP 詳解
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, "SVP — Shortest Vector Problem 最短向量問題", "Section 2 · Slide 14 / 41")
tb(s, "定義 / Definition", 0.45, 1.35, 12.4, 0.5,
   size=18, bold=True, color=NAVY)
tb(s, "給定一個格 Λ(B) ⊂ ℝⁿ，求非零向量 x ∈ Λ(B) \ {0} 使長度（歐氏範數）達到最小：",
   0.45, 1.85, 12.4, 0.5, size=16, color=BLACK)
formula_card(s, 0.9, 2.4, 11.5, 1.0,
             "λ₁(Λ)  =  min  ‖x‖    s.t.  x ∈ Λ \ {0}", size=22, color=GREEN)
tb(s, "直觀理解 / Intuition", 0.45, 3.6, 12.4, 0.5,
   size=18, bold=True, color=NAVY)
items_svp = [
    "Λ 是 ℝⁿ 中由基底整數線性組合構成的離散點集；SVP 即「在這些點中找離原點最近的非零點」",
    "二維且基底「漂亮」（接近正交、長度相近）時，最短向量很容易看出來",
    "若基底「歪扭」（基向量夾角小、長度差異大），最短向量仍存在卻難以辨識",
    "→ 此時需要 lattice basis reduction（如 LLL 演算法）先把基底整理成 good basis",
]
for i, t in enumerate(items_svp):
    y = 4.15 + i * 0.55
    tb(s, "▸", 0.9, y, 0.3, 0.4, size=15, bold=True, color=BLUE)
    tb(s, t, 1.2, y, 11.2, 0.5, size=14, color=BLACK)
tb(s, "💡  SVP 與其近似版本 SVP_γ 被認為對「古典 + 量子」都需 2^Θ(n) 時間。",
   0.55, 6.55, 12.2, 0.5, size=13, color=NAVY)

# Slide 15 — CVP 詳解
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, "CVP — Closest Vector Problem 最近向量問題", "Section 2 · Slide 15 / 41")
tb(s, "定義 / Definition", 0.45, 1.35, 12.4, 0.5,
   size=18, bold=True, color=NAVY)
tb(s, "給定格 Λ(B) ⊂ ℝⁿ 與目標點 y ∈ ℝⁿ（y 未必落在格上），求格點 x ∈ Λ(B) 使距離最小：",
   0.45, 1.85, 12.4, 0.5, size=15, color=BLACK)
formula_card(s, 0.9, 2.4, 11.5, 1.0,
             "dist(y, Λ)  =  min  ‖x − y‖    s.t.  x ∈ Λ(B)", size=22, color=GREEN)
tb(s, "直觀理解 / Intuition", 0.45, 3.6, 12.4, 0.5,
   size=18, bold=True, color=NAVY)
items_cvp = [
    "想像 ℝⁿ 鋪滿一格一格的離散晶格點，再任意丟下一點 y（多半不會落在格上）",
    "CVP 即「在這些晶格點中找離 y 最近的那一個」",
    "基底漂亮時：B⁻¹y 對每分量四捨五入，再乘回 B，即近似解",
    "基底歪扭時：四捨五入會選錯胞元（cell），需先做 lattice basis reduction",
    "SVP 與 CVP 的關係：CVP 至少和 SVP 一樣難（SVP 可規約至 CVP）",
]
for i, t in enumerate(items_cvp):
    y = 4.15 + i * 0.45
    tb(s, "▸", 0.9, y, 0.3, 0.4, size=15, bold=True, color=BLUE)
    tb(s, t, 1.2, y, 11.2, 0.4, size=14, color=BLACK)
tb(s, "💡  CVP 的困難性 → BDD → LWE：是後量子密碼 (Kyber、Dilithium) 的安全基礎。",
   0.55, 6.55, 12.2, 0.5, size=13, color=NAVY)

# Slide 16 — Minkowski's first theorem
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, "Minkowski 第一定理 (Minkowski's First Theorem)", "Section 2 · Slide 16 / 41")
tb(s, "幾何數論的奠基定理", 0.45, 1.4, 12.4, 0.5,
   size=18, bold=True, color=NAVY)
rect(s, 0.9, 2.0, 11.5, 1.5, LGRN)
rect(s, 0.9, 2.0, 0.08, 1.5, GREEN)
tb(s, "Theorem (Minkowski, 1889).", 1.1, 2.08, 11.2, 0.4,
   size=18, bold=True, color=GREEN, font=EN)
tb(s, "若  S ⊂ ℝⁿ  為中心對稱凸體，且  vol(S) > 2ⁿ · det(𝓛)，",
   1.1, 2.5, 11.2, 0.45, size=17, color=BLACK)
tb(s, "則  S  含有非零格點  v ∈ 𝓛 ∩ S \\ {0}。",
   1.1, 2.95, 11.2, 0.45, size=17, color=BLACK)
tb(s, "三個關鍵詞 / Three keywords：", 0.45, 3.8, 12.4, 0.45,
   size=17, bold=True, color=NAVY)
items = [
    ("中心對稱", "centrally symmetric:  x ∈ S  ⇒  −x ∈ S"),
    ("凸體", "convex body: 任兩點線段含於 S；S 為有界閉集"),
    ("體積足夠大", "vol(S) > 2ⁿ · det(𝓛)  為臨界條件"),
]
for i, (k, v) in enumerate(items):
    y = 4.35 + i * 0.5
    rect(s, 0.9, y, 2.6, 0.4, LBLUE)
    tb(s, k, 1.05, y + 0.04, 2.3, 0.35, size=15, bold=True, color=NAVY)
    tb(s, v, 3.7, y + 0.04, 8.7, 0.4, size=14, color=BLACK)
tb(s, "💡  下一頁將套用此定理至中心對稱立方體，推得 SVP 上界。",
   0.55, 6.55, 12.2, 0.5, size=14, color=NAVY)


# Slide 17 — Minkowski 視覺化
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, "Minkowski 第一定理 — 視覺化 (n = 2)", "Section 2 · Slide 17 / 41")
tb(s, "整數格 Λ = ℤ² 與中心對稱凸體（圓盤）:臨界體積 = 2² det(Λ) = 4",
   0.45, 1.32, 12.4, 0.45, size=16, bold=True, color=NAVY,
   align=PP_ALIGN.CENTER)
s.shapes.add_picture(
    r"D:\成大碩士班\114-2\數論（一）\期末報告\minkowski_viz.png",
    Inches(1.665), Inches(1.85), width=Inches(10))
rect(s, 0.4, 6.6, 12.5, 0.42, LBLUE)
tb(s, "💡  (a) vol(S) ≈ 2.27 ≤ 4:S 可能只含原點。"
      "  (b) vol(S) ≈ 6.16 > 4:S 必含非零格點 v ∈ Λ \\ {0}。",
   0.55, 6.62, 12.2, 0.4, size=13, color=NAVY)

# Slide 13 — SVP upper bound derivation Step 1
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, "推導 SVP 上界 — Step 1：選立方體", "Section 2 · Slide 18 / 41")
tb(s, "Apply Minkowski to a centered cube / 套用 Minkowski 至中心對稱立方體",
   0.45, 1.35, 12.4, 0.5, size=17, bold=True, color=NAVY)
formula_card(s, 0.9, 1.95, 11.5, 1.1,
             "Cᵣ := [-r, r]ⁿ ,    vol(Cᵣ) = (2r)ⁿ", size=22)
tb(s, "套用 Minkowski 定理的條件:vol(Cᵣ) > 2ⁿ · det(𝓛)",
   0.45, 3.25, 12.4, 0.5, size=17, color=BLACK)
formula_card(s, 0.9, 3.85, 11.5, 1.1,
             "(2r)ⁿ > 2ⁿ det(𝓛)   ⟺   r > det(𝓛)^(1/n)", size=22, color=GREEN)
tb(s, "因此只要  r > det(𝓛)^(1/n)，立方體 Cᵣ 即含有非零格點 v ∈ 𝓛。",
   0.45, 5.15, 12.4, 0.5, size=17, color=BLACK)
rect(s, 0.9, 5.85, 11.5, 0.7, LBLUE)
tb(s, "此 v 滿足  ‖v‖∞ ≤ r  — 因為  v ∈ Cᵣ = [-r, r]ⁿ。",
   1.1, 5.95, 11.2, 0.5, size=16, bold=True, color=NAVY)

# Slide 14 — SVP upper bound derivation Step 2 + conclusion
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, "推導 SVP 上界 — Step 2：範數不等式 + 結論", "Section 2 · Slide 19 / 41")
tb(s, "標準範數不等式 / Standard norm inequality",
   0.45, 1.35, 12.4, 0.5, size=17, bold=True, color=NAVY)
formula_card(s, 0.9, 1.95, 11.5, 1.0,
             "‖v‖₂  ≤  √n  ·  ‖v‖∞    (∀ v ∈ ℝⁿ)", size=22)
tb(s, "結合 Step 1 的  ‖v‖∞ ≤ r：",
   0.45, 3.1, 12.4, 0.5, size=17, color=BLACK)
formula_card(s, 0.9, 3.7, 11.5, 1.0,
             "‖v‖₂  ≤  √n · r ,    valid for all  r > det(𝓛)^(1/n)", size=20)
tb(s, "讓  r ↘ det(𝓛)^(1/n)，得結論：",
   0.45, 4.85, 12.4, 0.5, size=17, color=BLACK)
rect(s, 0.9, 5.5, 11.5, 1.1, LGRN)
rect(s, 0.9, 5.5, 0.08, 1.1, GREEN)
tb(s, "λ₁(𝓛)  ≤  √n  ·  det(𝓛)^(1/n)",
   1.1, 5.65, 11.2, 0.85, size=28, bold=True, color=GREEN,
   font=EN, align=PP_ALIGN.CENTER)
tb(s, "💡  Hermite 常數 γₙ 可給出更緊的上界  λ₁² ≤ γₙ det(𝓛)^(2/n)，γₙ ≤ 2n / πe + o(n)。",
   0.55, 6.7, 12.2, 0.5, size=13, color=NAVY)

# Slide 15 — Section 2 summary
bullets("Section 2 小結", [
    (1, "格  𝓛 = ⊕ ℤ bᵢ ⊂ ℝⁿ — 離散加法子群"),
    (2, "覆積  det(𝓛)、相繼極小  λᵢ  為核心不變量"),
    (1, "三大困難問題：SVP / GapSVP / BDD"),
    (2, "目前已知最佳量子演算法仍需  2^Θ(n)  時間"),
    (1, "Minkowski 第一定理 → SVP 上界"),
    (2, "λ₁(𝓛)  ≤  √n · det(𝓛)^(1/n)"),
    (2, "幾何（covolume）與算術（最短向量）的橋樑 — Course Objective 5"),
    (1, "下一節：把這套幾何裝到「分圓域 + 數環」上"),
    (2, "理想 → 理想格 → Ring-LWE 的攻擊難度由此而來"),
], sub="Section 2 · Slide 20 / 41")

# ════════════════════════════════════════════════════
# Section 3 — Cyclotomic Fields and Ideal Lattices (10 min)
# ════════════════════════════════════════════════════
section_divider(3, "分圓域與理想格",
                "Cyclotomic Fields and Ideal Lattices",
                "10 分鐘 / minutes")

# Slide 17 — General n-th cyclotomic
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, "一般 n-th Cyclotomic Field", "Section 3 · Slide 22 / 41")
tb(s, "本原 n 次單位根 / Primitive n-th root of unity",
   0.45, 1.35, 12.4, 0.5, size=18, bold=True, color=NAVY)
formula_card(s, 0.9, 1.95, 11.5, 1.0,
             "ζₙ  =  e^(2πi/n) ,    ζₙⁿ = 1 ,    ζₙᵏ ≠ 1   for  1 ≤ k < n", size=20)
tb(s, "分圓多項式 / Cyclotomic polynomial", 0.45, 3.15, 12.4, 0.5,
   size=18, bold=True, color=NAVY)
formula_card(s, 0.9, 3.75, 11.5, 1.0,
             "Φₙ(x)  =  ∏  (x − ζₙᵏ)   ∈   ℤ[x]\n                    gcd(k,n)=1", size=18)
items = [
    "Monic（首一）",
    "Irreducible over ℚ（在 ℚ 上不可約）",
    "deg Φₙ = φ(n)（Euler totient）",
]
for i, t in enumerate(items):
    y = 5.0 + i * 0.45
    tb(s, "▸", 0.9, y, 0.3, 0.4, size=16, bold=True, color=BLUE)
    tb(s, t, 1.2, y, 11, 0.4, size=16, color=BLACK)
tb(s, "💡  分圓多項式 {Φₙ} 滿足   xⁿ − 1 = ∏_{d ∣ n} Φ_d(x)。",
   0.55, 6.55, 12.2, 0.5, size=14, color=NAVY)

# Slide 18 — Examples
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, "分圓多項式範例 / Examples", "Section 3 · Slide 23 / 41")
examples = [
    ("Φ₁(x)",  "x − 1",                       "φ(1) = 1"),
    ("Φ₂(x)",  "x + 1",                       "φ(2) = 1"),
    ("Φ₃(x)",  "x² + x + 1",                  "φ(3) = 2"),
    ("Φ₄(x)",  "x² + 1",                      "φ(4) = 2"),
    ("Φ₅(x)",  "x⁴ + x³ + x² + x + 1",        "φ(5) = 4"),
    ("Φ₆(x)",  "x² − x + 1",                  "φ(6) = 2"),
    ("Φ₈(x)",  "x⁴ + 1",                      "φ(8) = 4"),
    ("Φ₁₂(x)", "x⁴ − x² + 1",                 "φ(12) = 4"),
]
tb(s, "Φₙ", 1.0, 1.4, 1.5, 0.4, size=18, bold=True, color=NAVY, font=EN)
tb(s, "形式 / Form", 3.0, 1.4, 6, 0.4, size=18, bold=True, color=NAVY)
tb(s, "次數", 10, 1.4, 2, 0.4, size=18, bold=True, color=NAVY)
rect(s, 0.9, 1.85, 11.5, 0.04, NAVY)
for i, (k, v, d) in enumerate(examples):
    y = 2.0 + i * 0.55
    if i % 2 == 0: rect(s, 0.9, y, 11.5, 0.5, LGRAY)
    tb(s, k, 1.0, y + 0.05, 1.8, 0.4, size=17, bold=True, color=BLUE, font=EN)
    tb(s, v, 3.0, y + 0.05, 6.5, 0.4, size=16, color=BLACK, font=EN)
    tb(s, d, 10, y + 0.05, 2.3, 0.4, size=15, color=GRAY, font=EN)
tb(s, "💡  Power-of-two 子類 (n = 2ᵏ):Φ_{2ᵏ}(x) = x^(2^(k−1)) + 1 — 形式最簡，成為 Ring-LWE 實作首選。",
   0.55, 6.6, 12.2, 0.5, size=13, color=NAVY)

# Slide 19 — Cyclotomic field K
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, "分圓域 K = ℚ(ζₙ) 與整數環 𝒪_K",
       "Section 3 · Slide 24 / 41")
formula_card(s, 0.9, 1.4, 11.5, 1.0,
             "K  =  ℚ(ζₙ) ,    [K : ℚ]  =  φ(n)", size=22)
formula_card(s, 0.9, 2.55, 11.5, 1.0,
             "𝒪_K  =  ℤ[ζₙ]    (ring of integers,  Dedekind domain)", size=20)
tb(s, "Dedekind 性質 (Course Obj. 1, 2)", 0.45, 3.85, 12.4, 0.5,
   size=18, bold=True, color=NAVY)
items = [
    "𝒪_K 為 Noetherian、整封閉、Krull 維度 1",
    "每個非零理想可唯一分解為質理想之積:𝔞 = 𝔭₁^e₁ · 𝔭₂^e₂ · … · 𝔭ᵣ^eᵣ",
    "理想算術取代了元素算術 — 唯一分解性的「正確類比」",
]
for i, t in enumerate(items):
    y = 4.45 + i * 0.5
    tb(s, "▸", 0.9, y, 0.3, 0.4, size=16, bold=True, color=BLUE)
    tb(s, t, 1.2, y, 11, 0.4, size=16, color=BLACK)
tb(s, "💡  ℤ 在 ℚ(ζₙ) 中的「正確版本」就是 ℤ[ζₙ] — 因為 ζₙ 是首一整係數多項式 Φₙ 的根。",
   0.55, 6.55, 12.2, 0.5, size=13, color=NAVY)

# Slide 20 — Specialization to m = 2^k
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, "特殊化:power-of-two cyclotomic (m = 2ᵏ)",
       "Section 3 · Slide 25 / 41")
tb(s, "Ring-LWE 實作首選", 0.45, 1.35, 12.4, 0.5,
   size=20, bold=True, color=NAVY)
formula_card(s, 0.9, 1.95, 11.5, 1.0,
             "m  =  2ᵏ  (k ≥ 2) ,    Φₘ(x)  =  x^(2^(k−1)) + 1", size=22)
tb(s, "符號慣例 (LPR convention)", 0.45, 3.15, 12.4, 0.5,
   size=18, bold=True, color=NAVY)
formula_card(s, 0.9, 3.75, 11.5, 1.4,
             "m  =  分圓指標 (cyclotomic index)\n"
             "n  :=  φ(m)  =  2^(k−1)  =  域度 [K : ℚ]",
             size=18)
items = [
    ("k = 2", "m = 4",     "n = 2",   "Φ₄ = x² + 1   (K = ℚ(i))"),
    ("k = 3", "m = 8",     "n = 4",   "Φ₈ = x⁴ + 1"),
    ("k = 4", "m = 16",    "n = 8",   "Φ₁₆ = x⁸ + 1"),
    ("k = 10","m = 1024",  "n = 512", "Φ₁₀₂₄ = x⁵¹² + 1（實用尺寸）"),
]
for i, (a, b, c, d) in enumerate(items):
    y = 5.35 + i * 0.32
    tb(s, a, 0.9, y, 0.9, 0.3, size=14, color=BLUE, font=EN)
    tb(s, b, 1.9, y, 1.4, 0.3, size=14, color=BLUE, font=EN)
    tb(s, c, 3.4, y, 1.4, 0.3, size=14, color=BLUE, font=EN)
    tb(s, d, 4.9, y, 7.4, 0.3, size=13, color=BLACK)

# Slide 21 — Splitting of primes
bullets("質理想分裂 (Splitting of Primes)  ·  Course Obj. 3", [
    (1, "對 p ∤ m：分裂模式由 p mod m 決定"),
    (2, "f := order of p modulo m in (ℤ/m)*"),
    (2, "p · 𝒪_K = 𝔭₁ · 𝔭₂ · … · 𝔭_g  共  g = n / f  個質理想"),
    (2, "每個 𝔭ᵢ 之 inertia degree 為 f，ramification index 為 1"),
    (1, "對 p = 2（m = 2ᵏ 時唯一分歧質數）"),
    (2, "2 · 𝒪_K = (1 − ζₘ)ⁿ — 完全分歧 (totally ramified)"),
    (1, "Ring-LWE 的關鍵選擇:q prime, q ≡ 1 (mod m)"),
    (2, "⇒ f = 1，q 在 𝒪_K 中完全分裂"),
    (2, "⇒ R_q := 𝒪_K / q · 𝒪_K  ≅  ∏ᵢ₌₁ⁿ ℤ_q（CRT）"),
    (2, "此即 NTT 的代數來源 — 詳見 Section 4"),
], sub="Section 3 · Slide 26 / 41",
   note="這是「質理想分裂模式」直接決定 Ring-LWE 演算法效率的具體例證。")

# Slide 22 — Trace, Norm, Discriminant
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, "Trace, Norm, Discriminant  ·  Course Obj. 4",
       "Section 3 · Slide 27 / 41")
items_def = [
    ("Trace 跡",       "Tr_{K/ℚ}(α)  =  Σᵢ σᵢ(α)"),
    ("Norm 範數",      "N_{K/ℚ}(α)   =  Πᵢ σᵢ(α)"),
    ("Different 差",  "𝔡_{K/ℚ}⁻¹  =  { α ∈ K : Tr(α · 𝒪_K) ⊆ ℤ }"),
    ("Discriminant", "Δ_K  =  N_{K/ℚ}(𝔡_{K/ℚ})"),
]
for i, (k, v) in enumerate(items_def):
    y = 1.4 + i * 0.85
    rect(s, 0.9, y, 3.3, 0.7, LBLUE)
    tb(s, k, 1.05, y + 0.18, 3.0, 0.4, size=16, bold=True, color=NAVY)
    tb(s, v, 4.4, y + 0.15, 8.0, 0.5, size=15, color=BLACK)
rect(s, 0.9, 4.95, 11.5, 1.4, LGRN)
rect(s, 0.9, 4.95, 0.08, 1.4, GREEN)
tb(s, "Power-of-two cyclotomic (m = 2ᵏ) discriminant：",
   1.1, 5.05, 11.2, 0.4, size=16, bold=True, color=GREEN)
tb(s, "|Δ_K|  =  2^(n(k−1)) ,    √|Δ_K|  =  2^(n(k−1)/2)",
   1.1, 5.5, 11.2, 0.4, size=20, bold=True, color=NAVY,
   font=EN, align=PP_ALIGN.CENTER)
tb(s, "m = 8, n = 4 (k = 3)：|Δ_K| = 2^(4·2) = 256，√|Δ_K| = 16",
   1.1, 5.95, 11.2, 0.4, size=14, color=BLACK)
tb(s, "💡  判別式測量「整數環的擾動」；分歧質數恰是 Δ_K 的質因數。",
   0.55, 6.55, 12.2, 0.5, size=13, color=NAVY)

# Slide 23 — Canonical embedding
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, "Canonical (Minkowski) Embedding — 把代數翻譯為幾何",
       "Section 3 · Slide 28 / 41")
tb(s, "對分圓 (r₁ = 0, r₂ = n/2):", 0.45, 1.35, 12.4, 0.5,
   size=18, bold=True, color=NAVY)
formula_card(s, 0.9, 1.95, 11.5, 1.4,
             "σ : K  ↪  K_ℝ := K ⊗_ℚ ℝ  ≅  ℂ^(n/2)  ≅  ℝⁿ\n"
             "α  ↦  (√2 Re σ₁(α),  √2 Im σ₁(α),  …)",
             size=18)
tb(s, "理想格 (Ideal Lattice)", 0.45, 3.55, 12.4, 0.5,
   size=18, bold=True, color=NAVY)
tb(s, "每個非零理想 𝔞 ⊆ 𝒪_K 在 σ 下成為滿秩格：",
   0.45, 4.05, 12.4, 0.5, size=16, color=BLACK)
formula_card(s, 0.9, 4.65, 11.5, 1.1,
             "det( σ(𝔞) )  =  N(𝔞)  ·  √|Δ_K|", size=24, color=GREEN)
tb(s, "💡  關鍵翻譯：理想 (代數結構)  →  滿秩格 (幾何結構)；範數 ↔ 覆積。"
      "Ring-LWE 攻擊難度即建基於此：理想格上的 SVP 仍困難。",
   0.55, 6.05, 12.2, 0.95, size=13, color=NAVY)

# Slide 24 — Worked example m=8
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, "Worked Example:m = 8, n = 4",
       "Section 3 · Slide 29 / 41")
tb(s, "K = ℚ(ζ₈) — 最小的非平凡 power-of-two cyclotomic",
   0.45, 1.35, 12.4, 0.5, size=18, bold=True, color=NAVY)
items = [
    ("分圓多項式", "Φ₈(x) = x⁴ + 1"),
    ("域度",      "[K : ℚ] = φ(8) = 4 = n"),
    ("整數環",    "𝒪_K = ℤ[ζ₈]"),
    ("判別式",    "|Δ_K| = 2^(4·2) = 256，√|Δ_K| = 16"),
]
for i, (k, v) in enumerate(items):
    y = 1.95 + i * 0.55
    rect(s, 0.9, y, 2.8, 0.5, LBLUE)
    tb(s, k, 1.05, y + 0.07, 2.5, 0.4, size=15, bold=True, color=NAVY)
    tb(s, v, 3.9, y + 0.05, 8.5, 0.45, size=15, color=BLACK)
rect(s, 0.9, 4.4, 11.5, 1.9, LGRN)
rect(s, 0.9, 4.4, 0.08, 1.9, GREEN)
tb(s, "選擇 q = 17 — 因為 17 ≡ 1 (mod 8)",
   1.1, 4.5, 11.2, 0.4, size=17, bold=True, color=GREEN)
tb(s, "▸  f = order of 17 mod 8 = 1    ⇒    g = n / f = 4",
   1.1, 4.95, 11.2, 0.4, size=16, color=BLACK)
tb(s, "▸  17 · 𝒪_K  =  𝔮₁ · 𝔮₂ · 𝔮₃ · 𝔮₄ — 完全分裂為 4 個質理想",
   1.1, 5.4, 11.2, 0.4, size=16, color=BLACK)
tb(s, "▸  R₁₇ := 𝒪_K / 17 · 𝒪_K  ≅  ℤ₁₇⁴  — 此即 NTT 結構（詳見 Section 4）",
   1.1, 5.85, 11.2, 0.4, size=16, color=BLACK)
tb(s, "💡  此小例完整連結了 Obj. 1–4 + Course Objective 5（理想格的覆積）。",
   0.55, 6.55, 12.2, 0.5, size=14, color=NAVY)

# ════════════════════════════════════════════════════
# Section 4 — LWE and Ring-LWE (10 min)
# ════════════════════════════════════════════════════
section_divider(4, "LWE 與 Ring-LWE",
                "LWE and Ring-LWE",
                "10 分鐘 / minutes")

# Slide 26 — LWE intuition
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, "LWE 直觀解釋 — 帶雜訊的線性方程組求解",
       "Section 4 · Slide 31 / 41")
tb(s, "Learning With Errors（學習雜訊干擾的線性函數）",
   0.45, 1.35, 12.4, 0.5, size=18, bold=True, color=NAVY)
rect(s, 0.9, 1.95, 11.5, 1.5, LBLUE)
rect(s, 0.9, 1.95, 0.08, 1.5, BLUE)
tb(s, "給定樣本  (aᵢ, bᵢ)  滿足  bᵢ ≈ ⟨aᵢ, s⟩ (mod q)，",
   1.1, 2.03, 11.2, 0.45, size=17, color=BLACK)
tb(s, "其中每個 bᵢ 都加了一個小的隨機誤差 eᵢ，",
   1.1, 2.5, 11.2, 0.45, size=17, color=BLACK)
tb(s, "目標：還原秘密向量 s。",
   1.1, 2.95, 11.2, 0.45, size=17, bold=True, color=NAVY)
rect(s, 0.9, 3.7, 5.6, 2.6, LGRN)
rect(s, 0.9, 3.7, 0.08, 2.6, GREEN)
tb(s, "若 eᵢ = 0（無雜訊）：", 1.1, 3.8, 5.2, 0.4,
   size=15, bold=True, color=GREEN)
tb(s, "高斯消去法", 1.1, 4.25, 5.2, 0.4, size=18, bold=True, color=BLACK)
tb(s, "多項式時間可解", 1.1, 4.7, 5.2, 0.4, size=14, color=BLACK)
tb(s, "O(n³)，明確線性代數", 1.1, 5.1, 5.2, 0.4, size=13, color=GRAY)
rect(s, 6.8, 3.7, 5.6, 2.6, LRED)
rect(s, 6.8, 3.7, 0.08, 2.6, RED)
tb(s, "若加上小雜訊：", 7.0, 3.8, 5.2, 0.4,
   size=15, bold=True, color=RED)
tb(s, "問題變為指數困難", 7.0, 4.25, 5.2, 0.4,
   size=18, bold=True, color=BLACK)
tb(s, "目前無有效演算法", 7.0, 4.7, 5.2, 0.4, size=14, color=BLACK)
tb(s, "包含量子演算法亦無能", 7.0, 5.1, 5.2, 0.4, size=13, color=GRAY)
tb(s, "💡  雜訊的存在「打破了線性結構」— 這正是格密碼安全性的來源。",
   0.55, 6.55, 12.2, 0.5, size=14, color=NAVY)

# Slide 27 — LWE formal definition
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, "LWE 形式定義 (Regev, STOC 2005 / JACM 2009)",
       "Section 4 · Slide 32 / 41")
tb(s, "參數 / Parameters：", 0.45, 1.3, 12.4, 0.4,
   size=17, bold=True, color=NAVY)
items_p = [
    "n, q ∈ ℤ⁺ — 維度與模數",
    "χ — 誤差分布 (typically discrete Gaussian, width αq, α < 1/√n)",
    "s ∈ ℤ_q^n — 固定的秘密向量",
]
for i, t in enumerate(items_p):
    y = 1.7 + i * 0.4
    tb(s, "▸ " + t, 0.9, y, 12, 0.4, size=14, color=BLACK)
formula_card(s, 0.9, 3.0, 11.5, 1.1,
             "(a,  ⟨a, s⟩ + e  mod q)   ∈   ℤ_q^n × ℤ_q\n"
             "a ← U(ℤ_q^n) ,    e ← χ", size=17)
tb(s, "兩個變體 / Two variants：", 0.45, 4.25, 12.4, 0.4,
   size=17, bold=True, color=NAVY)
rect(s, 0.9, 4.65, 5.6, 1.4, LBLUE)
tb(s, "Search-LWE", 1.1, 4.7, 5.2, 0.4, size=16, bold=True, color=NAVY)
tb(s, "由樣本還原秘密 s", 1.1, 5.1, 5.2, 0.4, size=14, color=BLACK)
tb(s, "Recover the secret s", 1.1, 5.5, 5.2, 0.4, size=12, color=GRAY, font=EN)
rect(s, 6.8, 4.65, 5.6, 1.4, LBLUE)
tb(s, "Decision-LWE", 7.0, 4.7, 5.2, 0.4, size=16, bold=True, color=NAVY)
tb(s, "區分 LWE 樣本與 uniform", 7.0, 5.1, 5.2, 0.4, size=14, color=BLACK)
tb(s, "Distinguish from uniform", 7.0, 5.5, 5.2, 0.4, size=12, color=GRAY, font=EN)
tb(s, "💡  Decision ⟺ Search 多項式時間等價（standard reduction）。",
   0.55, 6.55, 12.2, 0.5, size=13, color=NAVY)

# Slide 28 — Regev's reduction
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, "Regev 量子歸約 — 安全性的源頭",
       "Section 4 · Slide 33 / 41")
tb(s, "Worst-case to Average-case Reduction",
   0.45, 1.3, 12.4, 0.4, size=18, bold=True, color=NAVY)
formula_card(s, 0.9, 1.85, 11.5, 1.1,
             "LWE_{n,q,χ}    ⟵    GapSVP_γ ,  SIVP_γ\n"
             "on arbitrary n-dimensional lattices,    γ = Õ(n/α)",
             size=17)
tb(s, "意義 / Meaning：", 0.45, 3.15, 12.4, 0.4,
   size=17, bold=True, color=NAVY)
items = [
    "破解隨機 LWE 實例至少和解任意格上 worst-case SVP 一樣困難",
    "「最壞情形」的困難性轉移至「平均情形」— 罕見且強的安全保證",
    "歸約是量子的（用到量子傅立葉）— 古典歸約仍是開放問題",
    "γ 為近似因子；越小越強，但 LWE 安全參數隨之變嚴",
]
for i, t in enumerate(items):
    y = 3.6 + i * 0.55
    tb(s, "▸", 0.9, y, 0.3, 0.4, size=15, bold=True, color=BLUE)
    tb(s, t, 1.2, y, 11.2, 0.5, size=15, color=BLACK)
rect(s, 0.4, 6.15, 12.5, 0.85, LRED)
tb(s, "🚧  缺點 / Drawback：", 0.55, 6.2, 12.2, 0.4,
   size=15, bold=True, color=RED)
tb(s, "公鑰矩陣  A ∈ ℤ_q^(n×m)  大小  O(n²)，乘法 O(n²) — Ring-LWE 即為解此問題而生。",
   0.55, 6.55, 12.2, 0.45, size=14, color=BLACK)

# Slide 29 — Ring-LWE: Lift to ring of algebraic integers
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, "Ring-LWE — 提升至代數整數環",
       "Section 4 · Slide 34 / 41")
tb(s, "Lyubashevsky–Peikert–Regev (EUROCRYPT 2010 / JACM 2013)",
   0.45, 1.3, 12.4, 0.4, size=15, color=GRAY, font=EN)
tb(s, "把 ℤⁿ 換成分圓整數環", 0.45, 1.7, 12.4, 0.5,
   size=18, bold=True, color=NAVY)
formula_card(s, 0.9, 2.25, 11.5, 1.4,
             "R  :=  𝒪_K  =  ℤ[x] / Φₘ(x)  =  ℤ[x] / (xⁿ + 1)\n"
             "R_q  :=  R / qR\n"
             "R^∨  :=  𝔡_{K/ℚ}⁻¹    (codifferent,  Obj. 4)",
             size=17)
tb(s, "Ring-LWE 樣本", 0.45, 3.85, 12.4, 0.4,
   size=17, bold=True, color=NAVY)
formula_card(s, 0.9, 4.25, 11.5, 1.0,
             "(a,  a · s + e  mod qR^∨)   ∈   R_q × R_q^∨\n"
             "a ← U(R_q) ,    e ← Gaussian on K_ℝ",
             size=16)
items = [
    "secret  s ∈ R_q^∨、sample 中的 a 與 a · s + e 皆為環元素",
    "對 power-of-two  m = 2ᵏ：R^∨ = (1/n) R 為純量倍 → simplified Ring-LWE",
    "歸約：Ring-LWE  ⟵  approx-SVP on ideal lattices in R（量子）",
]
for i, t in enumerate(items):
    y = 5.5 + i * 0.4
    tb(s, "▸", 0.9, y, 0.3, 0.4, size=14, bold=True, color=BLUE)
    tb(s, t, 1.2, y, 11.2, 0.4, size=13, color=BLACK)

# Slide 30 — Prime splitting → NTT
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, "質理想分裂的角色 — NTT 的代數來源",
       "Section 4 · Slide 35 / 41")
tb(s, "選擇 q prime 滿足 q ≡ 1 (mod m)",
   0.45, 1.3, 12.4, 0.5, size=18, bold=True, color=NAVY)
tb(s, "則 Φₘ(x) 在 ℤ_q 中有 n 個相異根：",
   0.45, 1.85, 12.4, 0.45, size=16, color=BLACK)
formula_card(s, 0.9, 2.4, 11.5, 1.0,
             "Φₘ(x)  ≡  ∏ᵢ₌₁ⁿ (x − ωᵢ)   (mod q)", size=20)
tb(s, "由中國剩餘定理 (CRT)：",
   0.45, 3.6, 12.4, 0.45, size=16, color=BLACK)
formula_card(s, 0.9, 4.05, 11.5, 1.0,
             "R_q  =  ℤ[x] / (Φₘ(x), q)   ≅   ∏ᵢ₌₁ⁿ ℤ_q",
             size=20, color=GREEN)
rect(s, 0.9, 5.25, 11.5, 1.1, LGRN)
rect(s, 0.9, 5.25, 0.08, 1.1, GREEN)
tb(s, "▸ 這個同構就是 Number-Theoretic Transform (NTT)！",
   1.1, 5.32, 11.2, 0.4, size=17, bold=True, color=GREEN)
tb(s, "▸ 計算複雜度 O(n log n)，類比 FFT 之於普通卷積",
   1.1, 5.7, 11.2, 0.4, size=15, color=BLACK)
tb(s, "▸ 幾何上：q · 𝒪_K 完全分裂為 n 個質理想（Obj. 3）",
   1.1, 6.05, 11.2, 0.4, size=15, color=BLACK)
tb(s, "💡  一個課程目標（Obj. 3）= 一個演算法加速。",
   0.55, 6.6, 12.2, 0.5, size=13, color=NAVY)

# Slide 31 — Negacyclic + key size
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, "Negacyclic 卷積 與 公鑰大小節省",
       "Section 4 · Slide 36 / 41")
tb(s, "因 Φₘ(x) = xⁿ + 1，在 R_q 中 xⁿ ≡ −1，",
   0.45, 1.3, 12.4, 0.5, size=17, color=BLACK)
formula_card(s, 0.9, 1.85, 11.5, 1.1,
             "(a ∗ b)_k   =   Σ_{i+j=k} aᵢbⱼ   −   Σ_{i+j=k+n} aᵢbⱼ ,    0 ≤ k < n",
             size=17)
tb(s, "比較表 / Comparison Table",
   0.45, 3.2, 12.4, 0.5, size=18, bold=True, color=NAVY)
rows = [
    ("項目",       "LWE",                "Ring-LWE",            "節省"),
    ("公鑰大小",   "O(n²)",              "O(n)",                "factor n"),
    ("乘法複雜度", "O(n²)",              "O(n log n) (NTT)",    "factor n / log n"),
    ("安全歸約",   "any lattice SVP",    "ideal lattice SVP",   "略弱（仍開放）"),
    ("結構",       "純線性代數",          "代數結構（環）",        "—"),
]
for i, row in enumerate(rows):
    y = 3.8 + i * 0.45
    if i == 0:
        rect(s, 0.9, y, 11.5, 0.4, NAVY)
        for j, c in enumerate(row):
            tb(s, c, 1.05 + j * 2.85, y + 0.05, 2.7, 0.3,
               size=14, bold=True, color=WHITE)
    else:
        if i % 2 == 0: rect(s, 0.9, y, 11.5, 0.4, LGRAY)
        for j, c in enumerate(row):
            tb(s, c, 1.05 + j * 2.85, y + 0.05, 2.7, 0.35,
               size=13, color=BLACK)
tb(s, "💡  Ring-LWE 用「結構」換「效率」— 但結構也可能被攻擊利用，這是仍在研究的議題。",
   0.55, 6.55, 12.2, 0.5, size=13, color=NAVY)

# Slide 32 — Section 4 summary
bullets("Section 4 小結", [
    (1, "LWE：帶雜訊的線性方程組求解"),
    (2, "Regev 量子歸約：worst-case SVP → average-case LWE"),
    (2, "缺點：公鑰  O(n²)"),
    (1, "Ring-LWE：把 ℤⁿ 換為分圓整數環 R = 𝒪_K"),
    (2, "歸約：approx-SVP on ideal lattices"),
    (2, "公鑰縮為 O(n)、乘法 O(n log n)（NTT）"),
    (1, "代數結構的具體貢獻"),
    (2, "質理想分裂 (Obj. 3) → CRT 同構 → NTT"),
    (2, "Codifferent (Obj. 4) → dual ring R^∨；power-of-two 時退化為純量"),
], sub="Section 4 · Slide 37 / 41")

# ════════════════════════════════════════════════════
# Section 5 — Applications (5 min)
# ════════════════════════════════════════════════════
section_divider(5, "應用",
                "Applications",
                "5 分鐘 / minutes")

# Slide 34 — Why this matters
bullets("為什麼這件事重要？(Why This Matters)", [
    (1, "今日加密生態完全依賴古典公鑰"),
    (2, "HTTPS（網站）、行動銀行、即時通訊、軟體更新、VPN — 全部依賴 RSA / ECC"),
    (2, "量子電腦上線後將「同時」失效"),
    (1, "Harvest-Now, Decrypt-Later 威脅"),
    (2, "今日攔截的密文，明日量子電腦解密"),
    (2, "醫療紀錄、外交電文、長期密鑰 — 已暴露"),
    (1, "汰換時程"),
    (2, "NIST 建議 2030 前完成關鍵系統 PQC 轉換"),
    (2, "Google、Cloudflare、Apple 已開始於 TLS 部署 ML-KEM 混合模式"),
], sub="Section 5 · Slide 39 / 41",
   note="「先攔截、後解密」對長期機密是最棘手的威脅。")

# Slide 35 — NIST FIPS 203 / 204
s = blank(); rect(s, 0, 0, 13.33, 7.5, WHITE)
header(s, "NIST 後量子標準 (2024 年 8 月公布)",
       "Section 5 · Slide 40 / 41")
tb(s, "歷時 8 年的全球競賽 (2016–2024)",
   0.45, 1.3, 12.4, 0.5, size=15, color=GRAY)
rect(s, 0.9, 1.85, 11.5, 2.1, LGRN)
rect(s, 0.9, 1.85, 0.08, 2.1, GREEN)
tb(s, "FIPS 203 — 金鑰交換 / Key Encapsulation",
   1.1, 1.95, 11.2, 0.4, size=18, bold=True, color=GREEN)
tb(s, "▸ 後量子版本的 RSA / Diffie–Hellman 握手",
   1.1, 2.4, 11.2, 0.4, size=15, color=BLACK)
tb(s, "▸ 用於兩方建立加密通道前的「秘密握手」",
   1.1, 2.8, 11.2, 0.4, size=15, color=BLACK)
tb(s, "▸ 數學基礎：Ring-LWE / Module-LWE",
   1.1, 3.2, 11.2, 0.4, size=15, color=BLACK)
tb(s, "▸ 公鑰約 800 byte，密文約 1 KB；速度快於 RSA",
   1.1, 3.6, 11.2, 0.4, size=14, color=GRAY)
rect(s, 0.9, 4.1, 11.5, 2.1, LBLUE)
rect(s, 0.9, 4.1, 0.08, 2.1, BLUE)
tb(s, "FIPS 204 — 數位簽章 / Digital Signature",
   1.1, 4.2, 11.2, 0.4, size=18, bold=True, color=BLUE)
tb(s, "▸ 後量子版本的 RSA / ECDSA 簽章",
   1.1, 4.65, 11.2, 0.4, size=15, color=BLACK)
tb(s, "▸ 用於 TLS 憑證、軟體簽章、身分認證",
   1.1, 5.05, 11.2, 0.4, size=15, color=BLACK)
tb(s, "▸ 數學基礎：Ring-LWE / Module-LWE",
   1.1, 5.45, 11.2, 0.4, size=15, color=BLACK)
tb(s, "▸ 簽章約 2.4 KB，公鑰約 1.3 KB",
   1.1, 5.85, 11.2, 0.4, size=14, color=GRAY)
tb(s, "💡  代數數論 + 格幾何 = 你我未來十年的網路安全基礎。",
   0.55, 6.55, 12.2, 0.5, size=14, bold=True, color=NAVY)

# Slide 36 — Closing
s = blank(); rect(s, 0, 0, 13.33, 7.5, NAVY)
rect(s, 0, 3.2, 13.33, 0.05, GOLD)
tb(s, "Thank you!", 0, 1.2, 13.33, 1.0,
   size=64, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=EN)
tb(s, "感謝聆聽 · Q & A", 0, 2.35, 13.33, 0.7,
   size=32, color=LBLUE, align=PP_ALIGN.CENTER)
tb(s, "本報告核心：以代數數論的觀點理解 Ring-LWE",
   0, 3.5, 13.33, 0.5, size=18, color=LBLUE, align=PP_ALIGN.CENTER)
tb(s, "格 → 理想格 → Ring-LWE → 後量子密碼學",
   0, 4.0, 13.33, 0.5, size=16, color=LBLUE, align=PP_ALIGN.CENTER)
tb(s, "主要參考文獻 / Main References",
   0.5, 4.85, 12.3, 0.4, size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
refs = [
    "O. Regev, J. ACM 56(6), 2009.",
    "V. Lyubashevsky, C. Peikert, O. Regev, J. ACM 60(6), 2013.",
    "C. Peikert, Found. Trends TCS 10(4), 2016.",
    "NIST FIPS 203 / 204, 2024.    ·    D. A. Marcus, Number Fields, Ch. 1–5.",
]
for i, r in enumerate(refs):
    tb(s, r, 0.5, 5.25 + i * 0.32, 12.3, 0.3,
       size=12, color=LBLUE, align=PP_ALIGN.CENTER, font=EN)
tb(s, "黃崇晉  ·  L16141149  ·  數論（一）  ·  2026 春季  ·  國立成功大學",
   0, 6.7, 13.33, 0.4, size=12, color=LBLUE, align=PP_ALIGN.CENTER)

# ─── Save ──────────────────────────────────────────
import os
out = os.path.join(os.path.dirname(__file__) or ".", "presentation_v2.pptx")
prs.save(out)
print(f"Saved: {out}")
print(f"Total slides: {len(prs.slides)}")
